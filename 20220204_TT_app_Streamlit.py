# -*- coding: utf-8 -*-
"""
Created on Sat Oct 30 14:10:56 2021

@author: user
"""

import altair as alt 
import pydeck as pdk 
import streamlit as st 

import cvxpy as cp 
import dataframe_image as dfi 
import datetime as dt 
import pandas as pd 
import itertools 
import math 
import matplotlib.pyplot as plt 
import matplotlib.image as mpimg 
import matplotlib.ticker as mtick 
import numpy as np 
from datetime import datetime, timedelta
import import_ipynb 
import pandas_market_calendars as mcal 
import random 
from scipy.optimize import linprog 
import scipy.optimize as spo 
import pickle 
import pptx 
from pptx import Presentation 
from pptx.util import Inches 
from pptx.util import Pt 
from pptx.dml.color import RGBColor 
from pptx.enum.text import PP_ALIGN 
import PySimpleGUI as sg 
from PIL import Image 

# # API key for EIKON 
# import logging.config 
# import eikon as ek 
# ek.set_log_level(logging.DEBUG) 
 
# ek.set_app_key('6efe093f6ffd4ab59cf7bdd2fcd0674cf60c8bb6') 


def calcPortfolioLiquidationValue(dfPortfolioDataFiltered, shortTermTaxRate, longTermTaxRate): 
    # Calculate tax due on capital gains if entire portfolio was liquidated 
    # Different rates applied for long term and short term gains 
    taxDue = shortTermTaxRate * dfPortfolioDataFiltered['Short-term capital gains'].sum() + longTermTaxRate * dfPortfolioDataFiltered['Long-term capital gains'].sum() 
    
    # Calculate before and after tax portfolio value 
    beforeTaxValue = (dfPortfolioDataFiltered['Last price'] * dfPortfolioDataFiltered['Number of shares']).sum() 
    afterTaxValue = beforeTaxValue + taxDue 
    
    # Output 
    dfOutput = pd.DataFrame(index = ['Before tax value', 'After tax value', 'Tax liability'], columns = ['USD value']) 
    
    dfOutput.loc['Before tax value', 'USD value'] = round(beforeTaxValue, 2) 
    dfOutput.loc['After tax value', 'USD value'] = round(afterTaxValue, 2) 
    dfOutput.loc['Tax liability', 'USD value'] = round(taxDue, 2) 
  
    return dfOutput 


def replaceSecurity(dfPortfolioDataFiltered, constituentPrices, shortTermTaxRate, longTermTaxRate, currentDate, toBeReplacedIsin, replacementIsin, lastPriceOfStockToBeIncluded): 
    subset = dfPortfolioDataFiltered.loc[dfPortfolioDataFiltered['ISIN'] == toBeReplacedIsin, :] 
    dfResultLiquidationValue = calcPortfolioLiquidationValue(subset, shortTermTaxRate, longTermTaxRate) 
    valueBeforeTax = dfResultLiquidationValue.loc['Before tax value', 'USD value'] 
    valueAfterTax = dfResultLiquidationValue.loc['After tax value', 'USD value'] 
    taxImpact = dfResultLiquidationValue.loc['Tax liability', 'USD value'] 
    
    # Reset the corresponding row(s) of dfPortfolioDataFiltered 
    # Shares of replacement stock to buy = valueBeforeTax / lastPriceOfStockToBeIncluded (round down)
    # Return revised dfPortfolioDataFiltered and taxImpact 

    # SHOULD WE USE valueBeforeTax IF THE IMPACT OF TAX IS POSITIVE, I.E. IT INCREASES THE PORTFOLIO VALUE? 
    noOfNewShares = np.floor(valueAfterTax / lastPriceOfStockToBeIncluded) 
    
    if replacementIsin in list(dfPortfolioDataFiltered['ISIN'].values): 
        indexTemp = np.where(dfPortfolioDataFiltered['ISIN'] == replacementIsin)[0] 
        lstReplacement = [dfPortfolioDataFiltered.loc[indexTemp, 'Security name'], dfPortfolioDataFiltered.loc[indexTemp, 'Security ticker'], replacementIsin, currentDate, noOfNewShares, lastPriceOfStockToBeIncluded, lastPriceOfStockToBeIncluded, dfPortfolioDataFiltered.loc[indexTemp, 'GICS sector'], lastPriceOfStockToBeIncluded, lastPriceOfStockToBeIncluded, 0.0, 0.0, 0.0] 
    else: 
        dfTemp = ek.get_data([replacementIsin], fields = ['CF_NAME', 'TR.RIC', 'TR.GICSSector'])[0] 
        lstReplacement = [dfTemp.loc[0, 'CF_NAME'], dfTemp.loc[0, 'RIC'], replacementIsin, currentDate, noOfNewShares, lastPriceOfStockToBeIncluded, lastPriceOfStockToBeIncluded, dfTemp.loc[0, 'GICS Sector Name'], lastPriceOfStockToBeIncluded, lastPriceOfStockToBeIncluded, 0.0, 0.0, 0.0] 
    
    indicesToBeRevised = np.where(dfPortfolioDataFiltered['ISIN'] == toBeReplacedIsin)[0] 
    dfPortfolioDataFiltered.iloc[indicesToBeRevised[0], :] = lstReplacement 
    dfPortfolioDataFiltered = dfPortfolioDataFiltered.drop(indicesToBeRevised[1:]).reset_index(drop = True) 
    return [dfPortfolioDataFiltered, taxImpact] 


def withdrawCash(dfPortfolioDataFiltered, amountToBeWithdrawn): 
    remainingAmount = amountToBeWithdrawn 
    
    dfPortfolioDataFilteredSortedByTaxPerDollar = dfPortfolioDataFiltered.sort_values(by = 'Tax per dollar value') 
    dfPortfolioDataFilteredWithdrawCash = dfPortfolioDataFilteredSortedByTaxPerDollar.copy() 
    
    dfTradesForWithdrawingCash = pd.DataFrame() 
    while remainingAmount > 1: 
        for eachEntry in dfPortfolioDataFilteredSortedByTaxPerDollar.index: 
            if dfPortfolioDataFilteredSortedByTaxPerDollar.loc[eachEntry, 'Position value'] <= remainingAmount: 
                if dfTradesForWithdrawingCash.empty: 
                    dfTradesForWithdrawingCash = pd.DataFrame(dfPortfolioDataFilteredSortedByTaxPerDollar.loc[eachEntry]).T 
                else: 
                    dfTradesForWithdrawingCash = dfTradesForWithdrawingCash.append(pd.DataFrame(dfPortfolioDataFilteredSortedByTaxPerDollar.loc[eachEntry]).T) 
                
                dfPortfolioDataFilteredWithdrawCash = dfPortfolioDataFilteredWithdrawCash.drop(eachEntry) 
                
                remainingAmount = remainingAmount - dfPortfolioDataFilteredSortedByTaxPerDollar.loc[eachEntry, 'Position value'] 
            else: 
                dfTradesForWithdrawingCash = dfTradesForWithdrawingCash.join(pd.DataFrame(dfPortfolioDataFilteredSortedByTaxPerDollar.loc[eachEntry])) 
                
                if dfTradesForWithdrawingCash.empty: 
                    dfTradesForWithdrawingCash = pd.DataFrame(dfPortfolioDataFilteredSortedByTaxPerDollar.loc[eachEntry]).T 
                else: 
                    dfTradesForWithdrawingCash = dfTradesForWithdrawingCash.append(pd.DataFrame(dfPortfolioDataFilteredSortedByTaxPerDollar.loc[eachEntry]).T) 
    
                dfTradesForWithdrawingCash.loc[eachEntry, 'Number of shares'] = round(remainingAmount / dfPortfolioDataFilteredSortedByTaxPerDollar.loc[eachEntry, 'Last price']) 
    
                dfPortfolioDataFilteredWithdrawCash.loc[eachEntry, 'Number of shares'] = max(0, dfPortfolioDataFilteredWithdrawCash.loc[eachEntry, 'Number of shares'] - round(remainingAmount / dfPortfolioDataFilteredSortedByTaxPerDollar.loc[eachEntry, 'Last price'])) 
                
                dfTradesForWithdrawingCash.loc[eachEntry, 'Tax liability'] = dfTradesForWithdrawingCash.loc[eachEntry, 'Number of shares'] * dfPortfolioDataFiltered.loc[eachEntry, 'Last price'] * dfPortfolioDataFiltered.loc[eachEntry, 'Tax per dollar value'] 
                
                remainingAmount = 0.0 
                
            if remainingAmount <= 1: 
                break 
    
    return dfPortfolioDataFilteredWithdrawCash, dfTradesForWithdrawingCash 


def contributeCash(dfPortfolioDataFiltered, currentDate, cashToBeContributed): 
    # a = current price of each individual security in the portfolio 
    # b = current number of shares of each individual security in the portfolio 
    # c = current price & number of shares of each individual security in the portfolio 
    # p = vector of the constituents' prices 
    # tickers = vector of the constituents' tickers (corresponding to each price) 
    # Round sharesToAdd down to the closest integer (generally, you cannot buy a fraction of a share) 
    # This guarantees that we incur at most the amount of cash we wish to add 
    dfA = dfPortfolioDataFiltered.groupby('ISIN')[['Last price']].mean() 
    dfB = dfPortfolioDataFiltered.groupby('ISIN')[['Number of shares']].sum() 
    dfC = pd.concat([dfA, dfB], axis = 1) 
    isins = np.array(dfC.index.to_list()) 
    serPrices = pd.Series(list(itertools.chain(*dfA.values)), list(dfA.index)) 
    serTotalShares = pd.Series(list(itertools.chain(*dfB.values)), list(dfB.index)) 
    serComponentsValues = serPrices * serTotalShares 
    serProportionAddedForUsd1 = serComponentsValues / np.sum(serComponentsValues)                       # Shares of each constituent to be liquidated for getting USD 1 from the portfolio 
    serSharesToBeAdded = np.floor(cashToBeContributed * serProportionAddedForUsd1 / serPrices) 
    
    # Add serSharesToBeAdded for each stock 
    # To do so, create another dataframe dfNew to record the transactions and then merge it with the earlier positions 
    dfNewAdditions = pd.DataFrame(columns = dfPortfolioDataFiltered.columns, index = range(len(dfPortfolioDataFiltered), len(dfPortfolioDataFiltered) + len(isins), 1)) 
    for i in range(0, len(isins), 1): 
        relevantIndex = dfPortfolioDataFiltered[dfPortfolioDataFiltered['ISIN'] == isins[i]].index[0] 
        lstReplacement = [dfPortfolioDataFiltered.loc[relevantIndex, 'Security name'], dfPortfolioDataFiltered.loc[relevantIndex, 'Security ticker'], dfPortfolioDataFiltered.loc[relevantIndex, 'ISIN'], currentDate, serSharesToBeAdded[dfPortfolioDataFiltered.loc[relevantIndex, 'ISIN']], dfPortfolioDataFiltered.loc[relevantIndex, 'Last price'], dfPortfolioDataFiltered.loc[relevantIndex, 'Last price'], dfPortfolioDataFiltered.loc[relevantIndex, 'GICS sector'], dfPortfolioDataFiltered.loc[relevantIndex, 'Last price'], dfPortfolioDataFiltered.loc[relevantIndex, 'Last price'], 0.0, 0.0, 0.0] 
        dfNewAdditions.iloc[i, :] = lstReplacement 
    
    dfPortfolioDataFiltered = pd.concat([dfPortfolioDataFiltered, dfNewAdditions], axis = 0) 
    
    return dfPortfolioDataFiltered 


def harvestLosses(dfPortfolioDataFiltered, shortTermTaxRate, longTermTaxRate):
    # The function needs three inputs from previous algorithm 
    # The input portfolio is dfPortfolioDataFiltered from algorithm 1 
    # The input shortTermTaxRate come from previous data where it is used to calculate loss-harvesting 
    # The input longTermTaxRate come from previous data where it is used to calculate loss-harvesting 

    # Create empty short- and long-term cumulative losses 
    shortTermLoss = 0 
    longTermLoss = 0 

    # Copy oringinal portfolio for updation purposes 
    dfPortfolioUpdated = dfPortfolioDataFiltered.copy() 
    
    # Calculating Short term total losses from portfolio
    for eachValue in dfPortfolioDataFiltered['Short-term capital gains'].values: 
        if eachValue < 0: 
            shortTermLoss = shortTermLoss + eachValue 
    # Calculating long term total loss from portfolio 
    for eachValue in dfPortfolioDataFiltered['Long-term capital gains'].values: 
        if eachValue < 0: 
            longTermLoss = longTermLoss + eachValue 

    # Calculating total loss-harvesting based on tax rate
    lossHarvesting = longTermLoss * longTermTaxRate + shortTermLoss * shortTermTaxRate 
    
    # REMOVING STOCKS FROM THE PORTFOLIO WHERE THE CAPITAL LOSSES HAVE BEEN HARVESTED 

    lstIndicesToBeRemoved = np.where(dfPortfolioUpdated['Short-term capital gains'] < 0)[0].tolist() + np.where(dfPortfolioUpdated['Long-term capital gains'] < 0)[0].tolist() 
    
    lstIndicesToBeRetained = [] 
    for eachIndex in dfPortfolioUpdated.index: 
        # Replace short-term losses stock's date and cost price to today's date and market price 
        if eachIndex not in lstIndicesToBeRemoved: 
            lstIndicesToBeRetained = lstIndicesToBeRetained + [eachIndex] 
    
    return lossHarvesting, dfPortfolioUpdated.loc[lstIndicesToBeRetained]


def findAReplacer(currentIsin, dfReplacementEtfReturns, noOfYearsToCalculateReplacer, noOfSelectedReplacementEtfs): 
    # findAReplacer finds a replacement security for a given currentIsin (isin of the security to be replaced)
    
    # dfReplacementEtfReturns is the dataframe for the returns of the candiate replacement ETFs 
    # get the currentTicker's price and return data from Yahoo Finance
    arrReplacementIsins = np.array(dfReplacementEtfReturns.columns.to_list()) 
    startDate = datetime.today() - timedelta(days = 365 * noOfYearsToCalculateReplacer) 
    endDate = datetime.today() - timedelta(days = 1) 
    dfPricesStockToBeReplaced, dfTemp = ek.get_data([currentIsin], fields = ['TR.PriceClose.date', 'TR.PriceClose'], parameters = {'Frq': 'D', 'SDate': startDate.strftime('%Y-%m-%d'), 'CALCMETHOD': 'CLOSE', 'EDate': endDate.strftime('%Y-%m-%d')}) 
    dfPricesStockToBeReplaced = dfPricesStockToBeReplaced[['Date', 'Price Close']] 
    dfPricesStockToBeReplaced['Date'] = pd.to_datetime(dfPricesStockToBeReplaced['Date']).dt.date 
    dfPricesStockToBeReplaced = dfPricesStockToBeReplaced.set_index('Date') 
    dfReturnsStockToBeReplaced = (dfPricesStockToBeReplaced / dfPricesStockToBeReplaced.shift(1) - 1).dropna() 
    dfReturnsStockToBeReplaced = dfReturnsStockToBeReplaced.rename(columns = {dfReturnsStockToBeReplaced.columns[0] : currentIsin}) 
    
    if currentIsin not in list(dfReplacementEtfReturns.columns): 
        dfReplacementEtfReturns = dfReplacementEtfReturns.join(dfReturnsStockToBeReplaced) 
    
    varStockToBeReplaced = dfReturnsStockToBeReplaced.var()[dfReturnsStockToBeReplaced.columns[0]] 
    dfReplacementEtfReturns = dfReplacementEtfReturns.dropna() 
    dfCovarReplacementEtfReturns = dfReplacementEtfReturns.cov() 
    
    # Compare each of the ETFs in dfReplacementEtfReturns with currentIsin 
    # Compute the beta and intercept for each comparison and store the results in betas and intercepts 
    # Compute beta and intercept using the formulae from simple regression (WITHOUT PERFORMING ACTUAL OLS) 
    serVarStockToBeReplaced = pd.Series(varStockToBeReplaced, index = dfCovarReplacementEtfReturns[currentIsin].index) 
    serBetas = dfCovarReplacementEtfReturns[currentIsin] / serVarStockToBeReplaced 
    serReplacementEtfMeanReturns = dfReplacementEtfReturns.mean(axis = 0) 
    serDummyReturnsStockToBeReplaced = pd.Series(serReplacementEtfMeanReturns[currentIsin], index = serReplacementEtfMeanReturns.index) 
    serIntercepts = serReplacementEtfMeanReturns - serBetas * serDummyReturnsStockToBeReplaced 
    
    dfReplacementResults = pd.DataFrame(serBetas) 
    dfReplacementResults = dfReplacementResults.rename(columns = {dfReplacementResults.columns[0] : 'Slope'}) 
    dfTemp = pd.DataFrame(serIntercepts) 
    dfTemp = dfTemp.rename(columns = {dfTemp.columns[0] : 'Intercept'}) 
    dfReplacementResults = dfReplacementResults.join(dfTemp) 
    
    if currentIsin in dfReplacementResults.index: 
        dfReplacementResults = dfReplacementResults[dfReplacementResults.index != currentIsin] 

    # Ideally, we would like to see beta = 1 and intercept = 0 (currentTicker and ETF have roughly = daily returns) 
    # Thus, we characterize the goodness of fit by summing the absolute deviation of beta from 1 and intercept from 0 
    # For each ETF, we compute this sum; the ETF with the lowest such sum will be chosen as the replacement ticker 
    dfReplacementResults['Distance'] = abs(dfReplacementResults['Slope'] - 1) + abs(dfReplacementResults['Intercept']) 
    dfReplacementResults['Rank'] = dfReplacementResults['Distance'].rank(axis = 0, ascending = True) 
    dfReplacementResults = dfReplacementResults[dfReplacementResults['Rank'] <= noOfSelectedReplacementEtfs] 
    dfReplacementResults = dfReplacementResults.sort_values(by = ['Rank'], ascending = True)[['Slope', 'Intercept']] 
    
    return dfReplacementResults.loc[dfReplacementResults.index[0]], dfReplacementResults 


# Short term and long term tax rate 
def unrelizedGainLoss(dfPortfolio, shortTermTaxRate, longTermTaxRate): 
    # Create empty output data frame 
    dfUnrealized = pd.DataFrame(np.zeros((5, 1))) 
    dfUnrealized.index = ['Short-term gains', 'Short-term losses', 'Long-term gains', 'Long-term losses', 'Liquidation cost'] 
    dfUnrealized.columns = ['Unrealized gains / losses'] 

    # Calculate gains, losses and cost 
    dfUnrealized.loc['Short-term gains', 'Unrealized gains / losses'] = dfPortfolio['Short-term capital gains'][dfPortfolio['Short-term capital gains'] > 0].sum() 
    dfUnrealized.loc['Short-term losses', 'Unrealized gains / losses'] = dfPortfolio['Short-term capital gains'][dfPortfolio['Short-term capital gains'] < 0].sum() 
    dfUnrealized.loc['Long-term gains', 'Unrealized gains / losses'] = dfPortfolio['Long-term capital gains'][dfPortfolio['Long-term capital gains'] > 0].sum() 
    dfUnrealized.loc['Short-term gains', 'Unrealized gains / losses'] = dfPortfolio['Long-term capital gains'][dfPortfolio['Long-term capital gains'] < 0].sum() 
    dfUnrealized.loc['Liquidation cost', 'Unrealized gains / losses'] = np.sum([np.sum(dfUnrealized.iloc[0 : 2, 0]) * shortTermTaxRate, np.sum(dfUnrealized.iloc[2 : 4, 0]) * longTermTaxRate]) 

    return dfUnrealized 


# Calculating Liquidation cost on the portfolio 
def portfolioLiquidationCost(dfPortfolio, dfUnrealized): 
    totalValue = (dfPortfolio['Number of shares'] * dfPortfolio['Last price']).sum() 
    costRatio = (dfUnrealized.loc['Liquidation cost', 'Unrealized gains / losses'] / totalValue) * 100 
    dfCostRatio = pd.DataFrame(np.zeros((1, 1))) 
    dfCostRatio.iloc[0, 0] = costRatio 
    dfCostRatio.columns = ['Percentage'] 
    dfCostRatio.index = ['Liquidation cost ratio of portfolio'] 
    
    return dfCostRatio 


def portfolioShortOffset(dfPortfolio, dfUnrealized): 
    # Create empty output data frmae
    dfShortOffset = pd.DataFrame(np.zeros((3, 2))) 

    # Rename dataframe 
    dfShortOffset.index = ['Short-term securities with losses', 'Short-term securities offsetting losses', 'Short-term sum'] 
    dfShortOffset.columns = ['Value', 'Gain / loss value'] 

    # Calculate gains and losses and cost
    dfShortOffset.loc['Short-term securities with losses', 'Value'] = (dfPortfolio['Short-term capital gains'][dfPortfolio['Short-term capital gains'] < 0]).sum() 
    dfShortOffset.loc['Short-term securities offsetting losses', 'Value'] = (dfPortfolio['Short-term capital gains'][dfPortfolio['Short-term capital gains'] > 0]).sum() 
    dfShortOffset.loc['Short-term sum', 'Value'] = np.sum(dfShortOffset.iloc[0 : 2, 0]) 
    dfShortOffset.loc['Short-term securities with losses', 'Gain / loss value'] = dfUnrealized.loc['Short-term losses', 'Unrealized gains / losses'] 
    dfShortOffset.loc['Short-term securities offsetting losses', 'Gain / loss value'] = dfUnrealized.loc['Short-term gains', 'Unrealized gains / losses'] 
    dfShortOffset.loc['Short-term sum', 'Gain / loss value'] = np.sum(dfShortOffset.iloc[0 : 2, 1]) 
    
    return dfShortOffset 


def portfolioShortOffsetTaxCredit(dfUnrealized, shortTermTaxRate): 
    # Short term securities tax credit 
    taxCredit = np.sum(dfUnrealized.iloc[0 : 2, 0]) * shortTermTaxRate 
    dfTaxCredit = pd.DataFrame(np.zeros((1, 1))) 
    dfTaxCredit.iloc[0,0] = taxCredit 
    dfTaxCredit.columns = ['Tax credit'] 
    dfTaxCredit.index = ['Short-term securities selling'] 
    
    return dfTaxCredit 


# Setting benchmark return function 
def benchmark(lstBenchmarkIsins, lstBenchmarkWeights, startDate, endDate): 
    # dfBenchmarkPricesDump, dfTemp = ek.get_data(lstBenchmarkIsins, fields = ['TR.PriceClose.date', 'TR.PriceClose'], parameters = {'Frq': 'D', 'SDate': startDate.strftime('%Y-%m-%d'), 'CALCMETHOD': 'CLOSE', 'EDate': endDate.strftime('%Y-%m-%d')}) 
    dfBenchmarkPricesDump = pd.read_excel('20220206 Benchmark prices.xlsx', header = 0, index_col = 0) 

    lstTempDates = list(dfBenchmarkPricesDump['Date'].unique()) 
    
    while '' in lstTempDates: 
            lstTempDates.remove('') 
    
    dfBenchmarkComponentsPrices = pd.DataFrame(index = lstTempDates, columns = lstBenchmarkIsins) 
    dfBenchmarkComponentsPrices = dfBenchmarkComponentsPrices.sort_index() 
    for eachIsin in dfBenchmarkComponentsPrices.columns: 
        dfTemp = dfBenchmarkPricesDump[dfBenchmarkPricesDump['Instrument'] == eachIsin][['Date', 'Price Close']].set_index('Date') 
        serTemp = dfTemp.rename(columns = {'Price Close' : eachIsin})[eachIsin] 
        serTemp = serTemp[~serTemp.index.duplicated(keep = 'first')] 
        dfBenchmarkComponentsPrices[eachIsin] = serTemp.dropna() 
    
    dfBenchmarkComponentsPrices['Date'] = dfBenchmarkComponentsPrices.index 
    
    dfBenchmarkComponentsPrices['Date'] = pd.to_datetime(dfBenchmarkComponentsPrices['Date']).dt.date 
    
    dfBenchmarkComponentsPrices = dfBenchmarkComponentsPrices.set_index('Date') 
    
    dfBenchmarkComponentsReturns = (dfBenchmarkComponentsPrices / dfBenchmarkComponentsPrices.shift(1) - 1).dropna() 
    
    dfBenchmarkWeights = pd.DataFrame(index = dfBenchmarkComponentsReturns.index, columns = dfBenchmarkComponentsReturns.columns) 
    
    i = 0 
    for eachComponent in dfBenchmarkWeights.columns: 
        dfBenchmarkWeights[eachComponent] = lstBenchmarkWeights[i] 
        i += 1 
    
    serBenchmarkReturns = (dfBenchmarkComponentsReturns * dfBenchmarkWeights).sum(axis = 1) 
    
    dfBenchmarkReturns = pd.DataFrame(serBenchmarkReturns, columns = ['Benchmark return']) 

    return dfBenchmarkReturns 
    

# Setting update portfolio return function 
def updatePortfolioReturn(lstIsins, startDate, endDate): 
    # dfPortfolioComponentsPricesDump, dfTemp = ek.get_data(lstIsins, fields = ['TR.PriceClose.date', 'TR.PriceClose'], parameters = {'Frq': 'D', 'SDate': startDate.strftime('%Y-%m-%d'), 'CALCMETHOD': 'CLOSE', 'EDate': endDate.strftime('%Y-%m-%d')}) 
    dfPortfolioComponentsPricesDump = pd.read_excel('20220206 Portfolio component prices.xlsx', header = 0, index_col = 0) 

    lstTempDates = list(dfPortfolioComponentsPricesDump['Date'].unique()) 
    
    while '' in lstTempDates: 
        lstTempDates.remove('') 
    
    dfPortfolioComponentsPrices = pd.DataFrame(index = lstTempDates, columns = lstIsins) 
    dfPortfolioComponentsPrices = dfPortfolioComponentsPrices.sort_index() 
    for eachIsin in dfPortfolioComponentsPrices.columns: 
        dfTemp = dfPortfolioComponentsPricesDump[dfPortfolioComponentsPricesDump['Instrument'] == eachIsin][['Date', 'Price Close']].set_index('Date') 
        serTemp = dfTemp.rename(columns = {'Price Close' : eachIsin})[eachIsin] 
        serTemp = serTemp[~serTemp.index.duplicated(keep = 'first')] 
        dfPortfolioComponentsPrices[eachIsin] = serTemp.dropna() 
    
    dfPortfolioComponentsPrices['Date'] = dfPortfolioComponentsPrices.index 
    
    dfPortfolioComponentsPrices['Date'] = pd.to_datetime(dfPortfolioComponentsPrices['Date']).dt.date 
    
    dfPortfolioComponentsPrices = dfPortfolioComponentsPrices.set_index('Date') 
    
    dfPortfolioComponentsReturns = (dfPortfolioComponentsPrices / dfPortfolioComponentsPrices.shift(1) - 1).dropna() 

    return dfPortfolioComponentsReturns 


# Setting portfolio tracking error function 
def portfolioTrackingError(dfPortfolioData, dfPortfolioComponentsReturns, dfBenchmarkReturns): 
    # Create portfolio weight based on updated portfolio 
    # serComponentsWeights = dfPortfolioData['Number of shares'] * dfPortfolioData['Last price'] / (dfPortfolioData['Number of shares'] * dfPortfolioData['Last price']).sum() 
    dfPortfolioData['Weights'] = dfPortfolioData['Number of shares'] * dfPortfolioData['Last price'] / (dfPortfolioData['Number of shares'] * dfPortfolioData['Last price']).sum() 
    dfConsolidatedWeights = dfPortfolioData[['ISIN', 'Weights']].copy() 
    serConsolidatedWeightsGrouped = dfConsolidatedWeights.groupby(by = 'ISIN')['Weights'].sum() 
    
    dfPortfolioComponentsReturnsFinal = pd.DataFrame() 
    for eachIsin in serConsolidatedWeightsGrouped.index: 
        currentColIndex = list(dfPortfolioComponentsReturns.columns).index(eachIsin) 
        
        if dfPortfolioComponentsReturnsFinal.empty: 
            dfPortfolioComponentsReturnsFinal = pd.DataFrame(dfPortfolioComponentsReturns.iloc[:, currentColIndex]) 
        else: 
            dfPortfolioComponentsReturnsFinal = dfPortfolioComponentsReturnsFinal.join(pd.DataFrame(dfPortfolioComponentsReturns.iloc[:, currentColIndex])) 
    
    dfPortfolioComponentsReturnsFinal = dfPortfolioComponentsReturnsFinal.loc[~dfPortfolioComponentsReturnsFinal.index.duplicated(), :] 
    dfBenchmarkReturns = dfBenchmarkReturns.loc[~dfBenchmarkReturns.index.duplicated(), :] 
    
    dfPortfolioBenchmarkDiff = dfPortfolioComponentsReturnsFinal.subtract(np.array(dfBenchmarkReturns), fill_value = 0) 
    dfPortfolioBenchmarkCovMatrix = dfPortfolioBenchmarkDiff.cov() 
    trackingErrorSquared = np.array(serConsolidatedWeightsGrouped) @ dfPortfolioBenchmarkCovMatrix @ np.array(serConsolidatedWeightsGrouped).T 
    trackingError = np.sqrt(trackingErrorSquared) * np.sqrt(252) 

    return trackingError 


def shortOffsetPortfolioTe(dfPortfolio, lstBenchmarkIsins, lstBenchmarkWeights, noOfYearsForTrackingErrorCalculation): 
    # The next 3 lines find the update portfolio after sell all short-term securities 
    arrRmNpIndex = np.array(np.where(dfPortfolio['Short-term capital gains'] != 0), ndmin = 0) 
    arrRmLsIndex = arrRmNpIndex.tolist()[0] 
    dfShortPortfolioData = dfPortfolio.drop(labels = arrRmLsIndex) 

    # Setting return start and ending time 
    startDate = datetime.today() - timedelta(365 * noOfYearsForTrackingErrorCalculation) 
    endDate = datetime.today() 

    # Find benchmark return and updated short-term portfolio return
    dfBenchmarkReturns = benchmark(lstBenchmarkIsins, lstBenchmarkWeights, startDate, endDate) 
    dfPortfolioComponentsReturns = updatePortfolioReturn(dfShortPortfolioData, startDate, endDate) 

    # Calculating Short term securities selling tracking error 
    shortLiquidationTe = portfolioTrackingError(dfShortPortfolioData, dfPortfolioComponentsReturns, dfBenchmarkReturns) 
    dfShortLiquidationTe = pd.DataFrame(shortLiquidationTe, index = ['Short-term securities liquidation'], columns = ['Tracking error']) 

    return dfShortLiquidationTe 


def portfolioShortLongOffset(dfPortfolio, dfUnrealized): 
    # Create empty output data frmae 
    dfShortLongOffset = pd.DataFrame(np.zeros((6, 2))) 

    # Rename dataframe 
    dfShortLongOffset.index = ['Short-term securities with losses', 'Short-term securities offsetting losses', 'Short-term sum', 'Long-term securities with losses', 'Long-term securities offsetting losses', 'Long-term sum'] 
    dfShortLongOffset.columns = ['Value', 'Gain / loss value'] 

    # Calculate gain, losses and cost 
    dfShortLongOffset.iloc[0, 0] = (dfPortfolio['Last price'][dfPortfolio['Short-term capital gains'] < 0] * dfPortfolio['Number of shares'][dfPortfolio['Short-term capital gains'] < 0]).sum() 
    dfShortLongOffset.iloc[1, 0] = (dfPortfolio['Last price'][dfPortfolio['Short-term capital gains'] > 0] * dfPortfolio['Number of shares'][dfPortfolio['Short-term capital gains'] > 0]).sum() 
    dfShortLongOffset.iloc[2, 0] = dfShortLongOffset.iloc[0 : 2, 0].sum() 
    dfShortLongOffset.iloc[3, 0] = (dfPortfolio['Last price'][dfPortfolio['Long-term capital gains'] < 0] * dfPortfolio['Number of shares'][dfPortfolio['Long-term capital gains'] < 0]).sum() 
    dfShortLongOffset.iloc[4, 0] = (dfPortfolio['Last price'][dfPortfolio['Long-term capital gains'] > 0] * dfPortfolio['Number of shares'][dfPortfolio['Long-term capital gains'] > 0]).sum() 
    dfShortLongOffset.iloc[5, 0] = np.sum(dfShortLongOffset.iloc[3 : 5, 0]) 
    dfShortLongOffset.iloc[0, 1] = dfUnrealized.iloc[1, 0] 
    dfShortLongOffset.iloc[1, 1] = dfUnrealized.iloc[0, 0] 
    dfShortLongOffset.iloc[2, 1] = np.sum(dfShortLongOffset.iloc[0 : 2, 1]) 
    dfShortLongOffset.iloc[4, 1] = dfUnrealized.iloc[2, 0] 
    dfShortLongOffset.iloc[5, 1] = np.sum(dfShortLongOffset.iloc[3 : 5, 1]) 

    return dfShortLongOffset 


def calcMinTe(dfPortfolioDataReturns, dfBenchmarkReturn, numberOfTeWeights, shortTermTaxRate, longTermTaxRate, dfPortfolio): 
    # Volatility = tracking error, tracking error is the transform of the volatility 

    # Set the covariance matrix input data 
    df = dfPortfolioDataReturns 
    # Computing difference between component returns and benchmark returns 
    df = dfPortfolioDataReturns.subtract(np.array(dfBenchmarkReturn), fill_value = 0) 
    # Compute covariance matrix (df being the dataframe of securities' historical returns) 
    covMatrix = np.array(df.cov()) 
    # Number of securities
    n = covMatrix.shape[0]
    # Number of variables to optimize (in this case we need to optimize the portfolio weights, thus n = number of securities in the portfolio) 
    w = cp.Variable(n) 

    # Objective 1: Portfolio volitility 
    # Set the volitality object: the volatility of portfolio that needs to be minimized 
    # w the variables we need to optimize, cov_matrix  is the variables cofficient which also is the input matrix 
    trackingError = cp.quad_form(w, covMatrix) 

    # Find the minimum of volitality (tracking error) 
    prob = cp.Problem( 
    cp.Minimize(trackingError), # minimize volatility (which also is the tracking error)
    # Constrains of the object 
    [cp.sum(w) == 1, w >= 0] # sum of weights variables = 1, each weight variable >= 0 
    ) 
    prob.solve()

    lstWeights = [float('%0.4f' % v) for v in w.value] # Set the value width, if value less than 1 basis point (ex: x < 0.0001), the value would be considered as 0 
    
    dfOptimizedWeights = pd.DataFrame(lstWeights) 
    dfOptimizedWeights.columns = ['Weight'] 
    dfOptimizedWeights = dfOptimizedWeights.join(pd.DataFrame(dfPortfolio['ISIN'])) 
    dfOptimizedWeights.columns = ['Weight', 'ISIN'] 

    # Optimal minimum volitility (tracking error) 
    minTeSqd = trackingError.value # Minmium value of the portfolio volitality (which also is the tracking error) 

    return np.sqrt(minTeSqd) * np.sqrt(250), dfOptimizedWeights 


def calcMinTaxPaymentForTargetWeights(dfCurrentPortfolio, dfTargetPortfolioWeights): 
    currentPortfolioValue = dfCurrentPortfolio['Position value'].sum() 
    
    dfTargetPortfolioUsd = dfTargetPortfolioWeights * currentPortfolioValue 
    dfTargetPortfolioUsd.columns = ['Position value'] 
    
    dfCurrentPortfolioUsd = pd.DataFrame(dfCurrentPortfolio.groupby(by = 'ISIN')['Position value'].sum()) 
    
    dfToBeTradedUsd = dfTargetPortfolioUsd - dfCurrentPortfolioUsd 
    
    dfToBeSold = dfToBeTradedUsd[dfToBeTradedUsd['Position value'] < 0] * (-1) 
    dfToBeBought = dfToBeTradedUsd[dfToBeTradedUsd['Position value'] >= 0] 
    
    dfPortfolioWhichRequiresSelling = pd.DataFrame() 
    dfTaxImpact = pd.DataFrame(0.0, index = list(dfPortfolioDataFiltered['Security ticker'].unique()), columns = ['USD value']) 
    for eachIsin in dfToBeSold.index: 
        dfPortfolioTemp = dfCurrentPortfolio[dfCurrentPortfolio['ISIN'] == eachIsin] 
        
        dfPortfolioTemp = dfPortfolioTemp.sort_values(by = ['Tax per dollar value'], ascending = True) 
        
        for eachEntry in dfPortfolioTemp.index: 
            if dfPortfolioTemp.loc[eachEntry, 'Position value'] >= dfToBeSold.loc[eachIsin, 'Position value']: 
                dfPortfolioTemp.loc[eachEntry, 'Number of shares'] = dfPortfolioTemp.loc[eachEntry, 'Number of shares'] - round(dfToBeSold.loc[eachIsin, 'Position value'] / dfPortfolioTemp.loc[eachEntry, 'Last price']) 
                
                dfTaxImpact.loc[dfPortfolioTemp.loc[eachEntry, 'Security ticker']] = dfTaxImpact.loc[dfPortfolioTemp.loc[eachEntry, 'Security ticker']] + round(dfToBeSold.loc[eachIsin, 'Position value'] / dfPortfolioTemp.loc[eachEntry, 'Last price']) * dfPortfolioTemp.loc[eachEntry, 'Last price'] * dfPortfolioTemp.loc[eachEntry, 'Tax per dollar value'] 
            
                dfToBeSold.loc[eachIsin, 'Position value'] = 0 
            else: 
                dfPortfolioTemp.loc[eachEntry, 'Number of shares'] = 0 
                
                dfTaxImpact.loc[dfPortfolioTemp.loc[eachEntry, 'Security ticker']] = dfTaxImpact.loc[dfPortfolioTemp.loc[eachEntry, 'Security ticker']] + dfPortfolioTemp.loc[eachEntry, 'Tax liability'] 
                
                dfToBeSold.loc[eachIsin, 'Position value'] = dfToBeSold.loc[eachIsin, 'Position value'] - dfPortfolioTemp.loc[eachEntry, 'Position value'] 
                
            if dfToBeSold.loc[eachIsin, 'Position value'] == 0: 
                break 
            
        if dfPortfolioWhichRequiresSelling.empty: 
            dfPortfolioWhichRequiresSelling = dfPortfolioTemp 
        else: 
            dfPortfolioWhichRequiresSelling = dfPortfolioWhichRequiresSelling.append(dfPortfolioTemp) 

    dfPortfolioWhichRequiresBuying = pd.DataFrame() 
    for eachIsin in dfToBeBought.index: 
        dfPortfolioTemp = dfCurrentPortfolio[dfCurrentPortfolio['ISIN'] == eachIsin].copy() 
        
        dfTemp = dfCurrentPortfolio[dfCurrentPortfolio['ISIN'] == eachIsin] 

        securityName = dfTemp.loc[dfTemp.index[0], 'Security name'] 

        securityTicker = dfTemp.loc[dfTemp.index[0], 'Security ticker'] 
        
        purchaseDate = datetime.today().strftime('%Y-%m-%d') 
        
        lastPrice = dfTemp.loc[dfTemp.index[0], 'Last price'] 
        
        gicsSector = dfTemp.loc[dfTemp.index[0], 'GICS sector'] 

        lstToBeBought = [securityName, securityTicker, eachIsin, purchaseDate, round(dfToBeBought.loc[eachIsin, 'Position value'] / lastPrice), lastPrice, lastPrice, gicsSector, lastPrice, lastPrice, dfCurrentPortfolio[dfCurrentPortfolio['ISIN'] == eachIsin]['Asset category'].iloc[0], 0.0, 0.0, 0.0, 0.0, 0.0, 0.0, 0.0] 
        
        serTemp = pd.Series(lstToBeBought, index = dfCurrentPortfolio.columns) 
        
        if dfPortfolioWhichRequiresBuying.empty: 
            dfPortfolioWhichRequiresBuying = dfPortfolioTemp.append(serTemp, ignore_index = True) 
        else: 
            dfPortfolioWhichRequiresBuying = dfPortfolioWhichRequiresBuying.append(dfPortfolioTemp).append(serTemp, ignore_index = True) 

    return dfPortfolioWhichRequiresSelling.append(dfPortfolioWhichRequiresBuying), dfTaxImpact 


def optimalTeVsTaxPayment(dfPortfolioDataReturns, numberOfTeWeights, shortTermTaxRate, longTermTaxRate, dfPortfolio): 
    dfPortfolio['Total capital gains tax'] = dfPortfolio['Long-term capital gains'] * longTermTaxRate + dfPortfolio['Short-term capital gains'] * shortTermTaxRate 

    dfPortfolio['Position value'] = dfPortfolio['Number of shares'] * dfPortfolio['Last price'] 
    
    dfPortfolio['Current weight'] = dfPortfolio['Position value'] / dfPortfolio['Position value'].sum() 
    
    arrCurrentWeights = np.array(dfPortfolio['Current weight']) 
    
    dfPortfolio['Tax per dollar'] = dfPortfolio['Total capital gains tax'] / dfPortfolio['Position value'] 
    
    arrTaxPerDollar = np.array(dfPortfolio['Tax per dollar']) 
    
    totalPortfolioValue = dfPortfolio['Position value'].sum() 
    
    # Set the covariance matrix imput data 
    df = dfPortfolioDataReturns 
    # compute covariance matrix (df being the dataframe of securities' historical returns) 
    covMatrix = np.array(df.cov()) 
    # Number of securities
    n = covMatrix.shape[0] 
    # All securities standard deviation 
    assetStds = np.sqrt(np.diag(covMatrix)) 
    # Number of variables to optimize (in this case we need to optimize the portfolio weights, thus n = number of securities in the portfolio) 
    w = cp.Variable(n) 

    # Objective: Portfolio tracking error 
    # Set the tracking error object --> portfolio tracking error that needs to be minimized 
    # w is the variable we need to optimize, cov_matrix is the variable cofficient which is also the input matrix 
    trackingErrorSqd = cp.quad_form(w, covMatrix) 

    # Find the minimum of volatility (tracking error) 
    taxLimit = cp.Parameter() # Maximum tax 
    prob = cp.Problem( 
        cp.Minimize(trackingErrorSqd), # Minimize volatility (which is also the tracking error)
        # Constraints of the object 
        [ cp.sum(w) == 1, 
          w >= 0, 
          cp.sum(cp.multiply(cp.pos(arrCurrentWeights - w), arrTaxPerDollar * totalPortfolioValue)) <= taxLimit ] # sum of weights variables = 1, each weight variable >= 0 
          # cp.sum(cp.pos(arrCurrentWeights - w) * totalPortfolioValue) <= taxLimit] # sum of weights variables = 1, each weight variable >= 0 
    ) 

    # Define a loop function that inputs a series of tax constraints and returns the crossponding optimal weight that minimizes tracking error 
    # The output will return 3 values, 1. the minimum tracking error, 2. the tax to be paid, and 3. the optimal weight 
    def solveForOptTe(taxPayment): 
        taxLimit.value = taxPayment 
        result = prob.solve() 
        return (np.sqrt(trackingErrorSqd.value) * np.sqrt(252), (sum((arrCurrentWeights - w.value) * np.where((arrCurrentWeights - w.value) > 0, 1, 0) * arrTaxPerDollar)) * totalPortfolioValue, w.value) 

    maxTaxPayment = dfPortfolio['Long-term capital gains'][dfPortfolio['Long-term capital gains'] > 0].sum() * longTermTaxRate + dfPortfolio['Short-term capital gains'][dfPortfolio['Short-term capital gains'] > 0].sum() * shortTermTaxRate 
    minTaxPayment = dfPortfolio['Long-term capital gains'][dfPortfolio['Long-term capital gains'] < 0].sum() * longTermTaxRate + dfPortfolio['Short-term capital gains'][dfPortfolio['Short-term capital gains'] < 0].sum() * shortTermTaxRate 

    arrTaxPaymentIntervals = np.linspace(minTaxPayment, maxTaxPayment, 100) # Divide the range between min and max of tax payment into 100 equal intervals 

    # Iterate the value to get the all result 
    i = 0 
    dfResults = pd.DataFrame(index = range(100), columns = ['Tracking error', 'Tax payment']) 
    dfWeights = pd.DataFrame(index = range(100), columns = dfPortfolio.index) 
    for eachTaxPayment in arrTaxPaymentIntervals: 
        tplOutput = solveForOptTe(eachTaxPayment) 
        dfResults.loc[i, 'Tracking error'] = tplOutput[0] 
        dfResults.loc[i, 'Tax payment'] = tplOutput[1] 
        dfWeights.loc[i] = tplOutput[2] 
        i = i + 1 

    return dfResults, dfWeights 


def portfolioTeTpPlot(outputDf): 
    # Plot our residuals against the actual values of hyg 
    fig, ax = plt.subplots(figsize = (16, 9), dpi = 100) 
    plt.scatter(outputDf['Tracking error'], outputDf['Tax payment'], color = 'g', marker = 11) 

    title = "Portfolio tax-payment vs tracking error" 
    plt.title(title) 
    plt.xlabel("Tracking error") 
    plt.ylabel("Tax payment") 
    ax.grid(b = True, color = 'black', linestyle = '-.', linewidth = 0.5, alpha = 0.2) 

    # Saving graph 
    plt.show() 
    
    
def calcPortfolioAfterTaxPayment(dfPortfolioDataFiltered, taxValue): 
    dfPortfolioDataFiltered = dfPortfolioDataFiltered.sort_values(by = ['Tax per dollar value'], ascending = True) 
    
    dfPortfolioAfterTaxPayment = dfPortfolioDataFiltered.copy() 
    dfTransactionsRequiredForTaxPayment = pd.DataFrame() 
    
    i = 0 
    runningTaxValue = taxValue 
    while 1: 
        currentIndex = dfPortfolioDataFiltered.index[i] 
        
        if dfPortfolioDataFiltered.loc[currentIndex, 'Tax liability'] <= runningTaxValue: 
            dfPortfolioAfterTaxPayment = dfPortfolioAfterTaxPayment.drop(currentIndex, axis = 0) 
            
            dfTransactionsRequiredForTaxPayment = dfTransactionsRequiredForTaxPayment.append(dfPortfolioDataFiltered.loc[currentIndex], ignore_index = True) 
            
            runningTaxValue = runningTaxValue - dfPortfolioDataFiltered.loc[currentIndex, 'Tax liability'] 
        else: 
            taxLiabilityPerShare = dfPortfolioDataFiltered.loc[currentIndex, 'Tax liability'] / dfPortfolioDataFiltered.loc[currentIndex, 'Number of shares'] 
            
            sharesToBeSold = round(runningTaxValue / taxLiabilityPerShare) 
            originalShares = dfPortfolioAfterTaxPayment.loc[currentIndex, 'Number of shares'] 
            
            lstToBeModified = ['Number of shares', 'Short-term capital gains', 'Long-term capital gains', 'Position value', 'Tax liability'] 
            for eachField in lstToBeModified: 
                dfPortfolioAfterTaxPayment.loc[currentIndex, eachField] = ((originalShares - sharesToBeSold) / originalShares) * dfPortfolioAfterTaxPayment.loc[currentIndex, eachField] 
            
            dfTransactionsRequiredForTaxPayment = dfTransactionsRequiredForTaxPayment.append(dfPortfolioDataFiltered.loc[currentIndex], ignore_index = True) 
            for eachField in lstToBeModified: 
                dfTransactionsRequiredForTaxPayment.loc[dfTransactionsRequiredForTaxPayment.index[-1], eachField] = sharesToBeSold / originalShares * dfTransactionsRequiredForTaxPayment.loc[dfTransactionsRequiredForTaxPayment.index[-1], eachField] 
            
            runningTaxValue = 0 
        
        if abs(runningTaxValue) <= 0.1: 
            break 
        
        i = i + 1 
            
    return dfPortfolioAfterTaxPayment, dfTransactionsRequiredForTaxPayment 


def changeDataframeFormat(dfInput, lstFormats): 
    i = 0 
    dfOutput = pd.DataFrame(index = dfInput.index, columns = dfInput.columns) 
    for eachCol in dfInput.columns: 
        for eachRow in dfInput.index: 
            if lstFormats[i] == '': 
                dfOutput.loc[eachRow, eachCol] = dfInput.loc[eachRow, eachCol] 
            else: 
                dfOutput.loc[eachRow, eachCol] = lstFormats[i].format(dfInput.loc[eachRow, eachCol]) 
        
        i = i + 1 
    
    return dfOutput 

    
#--------------------------------------------------------------
# CREATING GUI 
#--------------------------------------------------------------

filePath = 'C:/Users/user/Desktop/Work/Yayati/Tax transition app' 
fileNameOverall = 'all_data_op.txt' 
fileNameTransactions = 'inv_transactions_Plaid.xlsx' 
with open(filePath + '/' + fileNameOverall, 'rb') as myFile: 
    dataInput = pickle.load(myFile) 

dfPortfolioData = pd.read_excel(filePath + '/' + fileNameTransactions, header = 0) 

# dfPortfolioData = pd.read_excel('Aperio tax transition - revised.xlsx') 

# dfPortfolioDataFiltered = dfPortfolioData.loc[:, ['Security name', 'Security ticker', 'ISIN', 'Purchase date', 'Number of shares', 'Cost basis per share']] 

# Missing columns = 'ISIN' 
dfPortfolioDataFiltered = dfPortfolioData.loc[:, ['ticker_symbol', 'ticker_symbol', 'Quantity', 'Price', 'Date']] 
dfPortfolioDataFiltered.columns = ['Security name', 'Security ticker', 'Number of shares', 'Cost basis per share', 'Purchase date'] 

dfBenchmarkDetails = dataInput['client_benchmark'] 
count = len(dfBenchmarkDetails) 
dfClientBenchmark = pd.DataFrame(np.zeros((count, 2))) 
dfClientBenchmark.columns = ['Tickers', 'Weights'] 
i = 0 
for eachIndex in dfBenchmarkDetails.index: 
    dfClientBenchmark.loc[i, 'Tickers'] = dfBenchmarkDetails.loc[eachIndex, 'tickers'] 
    dfClientBenchmark.loc[i, 'Weights'] = float(dfBenchmarkDetails.loc[eachIndex, 'weights']) / 100 if float(dfBenchmarkDetails.loc[eachIndex, 'weights']) > 1.0 else float(dfBenchmarkDetails.loc[eachIndex, 'weights']) 
    
    i = i + 1 

lstPortfolioTickers = [tickerName + '@ticker' for tickerName in dfPortfolioDataFiltered['Security ticker'].values] 
lstBenchmarkTickers = [tickerName + '@ticker' for tickerName in dfClientBenchmark['Tickers'].values] 
# dfRics, dfTemp = ek.get_data(lstPortfolioTickers + lstBenchmarkTickers, ['TR.ISIN']) 
dfRics = pd.read_excel('20220206 RICS data.xlsx', header = 0, index_col = 0) 
# dfDataDump, dfTemp = ek.get_data(list(dfRics['ISIN']), ['TR.RIC', 'TR.ISIN', 'TR.PriceClose', 'TR.GICSSector', 'TR.AssetCategory', 'TR.InstrumentType']) 
dfDataDump = pd.read_excel('20220206 RIC ISIN Sector data.xlsx', header = 0, index_col = 0) 

dfPortfolioDataFiltered['ISIN'] = '' 
for eachIndex in dfPortfolioDataFiltered.index: 
    dfPortfolioDataFiltered.loc[eachIndex, 'ISIN'] = dfRics[dfRics['Instrument'] == dfPortfolioDataFiltered.loc[eachIndex, 'Security name']]['ISIN'].iloc[0] 

dfPortfolioDataFiltered['Refinitiv RIC'] = '' 
for eachIndex in dfPortfolioDataFiltered.index: 
    dfPortfolioDataFiltered.loc[eachIndex, 'Refinitiv RIC'] = dfDataDump[dfDataDump['ISIN'] == dfPortfolioDataFiltered.loc[eachIndex, 'ISIN']].loc[eachIndex, 'RIC'] 
    
dfPortfolioDataFiltered = dfPortfolioDataFiltered.drop(columns = ['Security ticker']) 

dfPortfolioDataFiltered = dfPortfolioDataFiltered.rename(columns = {'Refinitiv RIC' : 'Security ticker'}) 

if 'Cost basis per share' not in dfPortfolioDataFiltered.columns: 
    dfPortfolioDataFiltered['Cost basis per share'] = '' 

dfPortfolioDataFiltered = dfPortfolioDataFiltered[['Security name', 'Security ticker', 'ISIN', 'Purchase date', 'Number of shares', 'Cost basis per share']] 

allRics = pd.read_excel(open('RICs and ISINs.xlsx', 'rb')) 

# # Creating GUI 

# sg.theme("LightGreen5") 
# ft = ("Calibri", 9) 

fileFormat = ['Fidelity', 'Schwab', 'Others'] 

lstRicReplacement = list(allRics['RIC']) 
lstRicReplacement.sort() 
# dfPortfolioDataFilteredGroupedByIsin = dfPortfolioDataFiltered.groupby(by = 'ISIN')['Security ticker', 'Last price'].first() 
lstCurrentHoldings = list(dfPortfolioDataFiltered['Security ticker']) 
lstCurrentHoldings.sort() 

# # Benchmark customisation
# def new_layout(i): 
#     return [[sg.T("Benchmark: ", font = ft), sg.Combo(lstRicReplacement, key=("-bm-", i)), sg.T("Weight: (%)", font = ft), sg.InputText(0.0, key=("-wt-", i),justification='right')]] 

# columnLayout =  [
#                     [sg.T("Benchmark: ", font = ft), sg.Combo(lstRicReplacement, key=("-bm-", 0)), sg.T("Weight: (%)", font = ft), sg.InputText(0.0, key=("-wt-", 0),justification='right'), sg.Button("Add",enable_events=True,  key="-plus-")]
#                 ] 
# w, h = sg.Window.get_screen_size() 

# layout = [ 
#          [ sg.Image('logo_resized.png', size=(200, 100)) ], 
#          [ 
#              sg.Text("Client name", size = (50, 2), font = ft), 
#              sg.Input('', key = "-client_name-", justification = 'right', font = ft) 
#          ], 
#          [ sg.Text("Proportion of tax to be paid / assets transitioned (%)", size = (50, 2), font = ft), sg.Input(0.00, key = "-selected_proportion-", justification = 'right', font = ft) ], 
#          [ sg.Text("Withdraw cash from portfolio ($)", size = (50, 2), font = ft), sg.Input(0.00, key = "-withdraw_cash-", justification = 'right', font = ft) ], 
#          [ 
#              sg.Text ("Holding to be replaced", size = (50, 2), font = ft), 
#              sg.Combo(lstCurrentHoldings, enable_events = True, key = "-holding_tobe_replaced-") 
#          ], 
#          [ 
#              sg.Text ("Replacement stock", size = (50, 2), font = ft), 
#              sg.Combo(lstRicReplacement, enable_events = True, key = "-repacement_stock-") 
#          ], 
#          [ sg.Button("Submit", font = ft, pad = (425, 40)) ]] 

# window = sg.Window(title = "Tax", layout = layout, size = (w, h)) 

# # Creating event loop that closes window on clicking 'submit' or manually closing the window
# # Adding exceptions to check for benchmark weights adding to 100%, choosing client, 
# # And entering valid benchmark ticker

# count = 1 

# while True: 
#     event, values = window.read() 
    
#     if event == "Submit": 
#         break 
#     elif event == sg.WIN_CLOSED: 
#         break 
            
# event, values = window.read() 
# window.close() 

# clientName = values['-client_name-'] 
# clientShortTaxrate = float(dataInput['client_short_taxrate']) / 100 if float(dataInput['client_short_taxrate']) > 1.0 else float(dataInput['client_short_taxrate']) 
# clientLongTaxrate = float(dataInput['client_long_taxrate']) / 100 if float(dataInput['client_long_taxrate']) > 1.0 else float(dataInput['client_long_taxrate']) 
# # clientFileFormat = values['-file_format-'] 
# clientCashWithdrawal = float(values['-withdraw_cash-']) 
# clientHoldingToBeReplaced = values['-holding_tobe_replaced-'] 
# clientReplacementStock = values['-repacement_stock-'] 
# selectedProportion = float(values['-selected_proportion-']) / 100.0 

st.title('Tax transition app') 

clientName = st.text_input('Enter client name: ', value = 'Abc') 
clientShortTaxrate = st.number_input('Short-term tax rate (%): ', 0, 100, 43) / 100 
clientLongTaxrate = st.number_input('Long-term tax rate (%): ', 0, 100, 20) / 100 
clientFileFormat = st.selectbox('Choose format for the client''s files: ', ('Fidelity', 'Schwab', 'Others')) 
clientCashWithdrawal = st.number_input('Cash to be withdrawn by client: ', 0, 1000000, 100000) 
clientHoldingToBeReplaced = st.selectbox('Holding to be removed: ', (lstCurrentHoldings), index = 0) 
clientReplacementStock = st.selectbox('To be replaced by: ', (lstRicReplacement), index = 776) 
selectedProportion = st.number_input('Maximum tax bill proportion (%): ', 0, 100, 40) / 100 

#--------------------------------------------------------------
# MAIN CODE 
#--------------------------------------------------------------

#--------------------------------------------------------------
# Parameters specification 
#--------------------------------------------------------------

dfRicsAndIsins = pd.read_excel('RICs and ISINs.xlsx') 
dfRicsAndIsins = dfRicsAndIsins.set_index('RIC') 

# Getting client's tax rates 
shortTermTaxRate = clientShortTaxrate 
longTermTaxRate = clientLongTaxrate 
amountToBeWithdrawn = clientCashWithdrawal 
tickerToBeReplaced = clientHoldingToBeReplaced 
replacementTicker = clientReplacementStock 

dictBenchmarkData = {} 
for eachIndex in dfClientBenchmark.index: 
    dictBenchmarkData[dfClientBenchmark.loc[eachIndex, 'Tickers']] = dfClientBenchmark.loc[eachIndex, 'Weights'] 

# dictBenchmarkData =     { 
#                             'SPY': 0.90, 
#                             'AGG': 0.10 
#                         } 

noOfYearsToCalculateReplacer = 2 
noOfSelectedReplacementEtfs = 5 
noOfYearsForTrackingErrorCalculation = 2 
numOfGroupsForCalculatingTransitionTaxLiabilities = 10 

# Ensuring tax rate is decimal and not percent
if shortTermTaxRate > 1: 
    shortTermTaxRate = shortTermTaxRate / 100 
    
if longTermTaxRate > 1: 
    longTermTaxRate = longTermTaxRate / 100 

dfPortfolioDataFiltered['Last price'] = np.nan 
dfPortfolioDataFiltered['GICS sector'] = '' 
dfPortfolioDataFiltered['Instrument type'] = '' 
dfPortfolioDataFiltered['Asset category'] = '' 
dfPortfolioDataFiltered['Instrument type'] = '' 
dfPortfolioDataFiltered['Purchase date prices'] = np.nan 
dfPortfolioDataFiltered['Final purchase price'] = np.nan 
dfTempPrices = pd.read_excel('20220206 Purchase date prices.xlsx') 
# Download prices over 1 month ending on the purchase date and use the last known price as the cost basis 
for eachIndex in dfPortfolioDataFiltered.index: 
    dfPortfolioDataFiltered.loc[eachIndex, 'Last price'] = dfDataDump[dfDataDump['ISIN'] == dfPortfolioDataFiltered.loc[eachIndex, 'ISIN']]['Price Close'].iloc[0] 
    dfPortfolioDataFiltered.loc[eachIndex, 'GICS sector'] = dfDataDump[dfDataDump['ISIN'] == dfPortfolioDataFiltered.loc[eachIndex, 'ISIN']]['GICS Sector Name'].iloc[0]  
    dfPortfolioDataFiltered.loc[eachIndex, 'Instrument type'] = dfDataDump[dfDataDump['ISIN'] == dfPortfolioDataFiltered.loc[eachIndex, 'ISIN']]['Instrument Type'].iloc[0] 

    if type(dfPortfolioDataFiltered.loc[eachIndex, 'Purchase date']) == str: 
        purchaseDate = dt.datetime.strptime(dfPortfolioDataFiltered.loc[eachIndex, 'Purchase date'], '%Y-%m-%d') 
    else: 
        purchaseDate = dfPortfolioDataFiltered.loc[eachIndex, 'Purchase date'] 

    # dfPortfolioDataFiltered.loc[eachIndex, 'Purchase date prices'] = ek.get_data(dfPortfolioDataFiltered.loc[eachIndex, 'ISIN'], fields = ['TR.PriceClose'], parameters = {'Frq': 'D', 'SDate': purchaseDate.strftime('%Y-%m-%d'), 'CALCMETHOD': 'CLOSE'})[0].loc[0, 'Price Close'] 
    dfPortfolioDataFiltered.loc[eachIndex, 'Purchase date prices'] = dfTempPrices.loc[eachIndex, 'Purchase date prices'] 
    dfPortfolioDataFiltered.loc[eachIndex, 'Final purchase price'] = dfPortfolioDataFiltered.loc[eachIndex, 'Purchase date prices'] 
    if dfPortfolioDataFiltered.loc[eachIndex, 'Instrument type'] == 'Ordinary Shares' or dfPortfolioDataFiltered.loc[eachIndex, 'Instrument type'] == 'Equity ETFs': 
        dfPortfolioDataFiltered.loc[eachIndex, 'Asset category'] = 'Equities' 
    elif dfPortfolioDataFiltered.loc[eachIndex, 'Instrument type'] == 'Bond ETFs': 
        dfPortfolioDataFiltered.loc[eachIndex, 'Asset category'] = 'Bonds' 

dfPortfolioDataFiltered = dfPortfolioDataFiltered.drop(['Instrument type'], axis = 1) 

# Calculating long-term and short-term capital gains 
dfPortfolioDataFiltered['Holding period (years)'] = np.nan 
dfPortfolioDataFiltered['Long-term capital gains'] = 0 
dfPortfolioDataFiltered['Short-term capital gains'] = 0 

for eachIndex in dfPortfolioDataFiltered.index: 
    if type(dfPortfolioDataFiltered.loc[eachIndex, 'Purchase date']) == str: 
        purchaseDate = dt.datetime.strptime(dfPortfolioDataFiltered.loc[eachIndex, 'Purchase date'], '%Y-%m-%d').date() 
    else: 
        purchaseDate = dfPortfolioDataFiltered.loc[eachIndex, 'Purchase date'] 

    dfPortfolioDataFiltered.loc[eachIndex, 'Holding period (years)'] = (datetime.today().date() - purchaseDate.date()).days / 365.0 
    if dfPortfolioDataFiltered.loc[eachIndex, 'Holding period (years)'] > 1: 
        dfPortfolioDataFiltered.loc[eachIndex, 'Long-term capital gains'] = (dfPortfolioDataFiltered.loc[eachIndex, 'Last price'] - dfPortfolioDataFiltered.loc[eachIndex, 'Final purchase price']) * dfPortfolioDataFiltered.loc[eachIndex, 'Number of shares'] 
    else: 
        dfPortfolioDataFiltered.loc[eachIndex, 'Short-term capital gains'] = (dfPortfolioDataFiltered.loc[eachIndex, 'Last price'] - dfPortfolioDataFiltered.loc[eachIndex,'Final purchase price']) * dfPortfolioDataFiltered.loc[eachIndex, 'Number of shares'] 

# Processing data for the client benchmark 

dfClientBenchmark['ISIN'] = '' 
dfClientBenchmark['RIC'] = '' 
dfClientBenchmark['Last price'] = np.nan 
dfClientBenchmark['GICS sector'] = '' 
dfClientBenchmark['Instrument type'] = '' 
dfClientBenchmark['Asset category'] = '' 
for eachIndex in dfClientBenchmark.index: 
    dfClientBenchmark.loc[eachIndex, 'ISIN'] = dfRics[dfRics['Instrument'] == dfClientBenchmark.loc[eachIndex, 'Tickers']]['ISIN'].iloc[0] 
    dfClientBenchmark.loc[eachIndex, 'RIC'] = dfDataDump[dfDataDump['ISIN'] == dfClientBenchmark.loc[eachIndex, 'ISIN']]['RIC'].iloc[0] 
    dfClientBenchmark.loc[eachIndex, 'Last price'] = dfDataDump[dfDataDump['ISIN'] == dfClientBenchmark.loc[eachIndex, 'ISIN']]['Price Close'].iloc[0] 
    dfClientBenchmark.loc[eachIndex, 'GICS sector'] = dfDataDump[dfDataDump['ISIN'] == dfClientBenchmark.loc[eachIndex, 'ISIN']]['GICS Sector Name'].iloc[0] 
    dfClientBenchmark.loc[eachIndex, 'Instrument type'] = dfDataDump[dfDataDump['ISIN'] == dfClientBenchmark.loc[eachIndex, 'ISIN']]['Instrument Type'].iloc[0] 
    if dfClientBenchmark.loc[eachIndex, 'Instrument type'] == 'Ordinary Shares' or dfClientBenchmark.loc[eachIndex, 'Instrument type'] == 'Equity ETFs': 
        dfClientBenchmark.loc[eachIndex, 'Asset category'] = 'Equities' 
    elif dfClientBenchmark.loc[eachIndex, 'Instrument type'] == 'Bond ETFs': 
        dfClientBenchmark.loc[eachIndex, 'Asset category'] = 'Bonds' 

# CALCULATING NUMBERS FOR THE 55ip 1-PAGER (SLIDES 4, 5 and 7) 

# (i) Portfolio value, total tax payable and tracking error 

dfPortfolioValueAndTotalTax = pd.DataFrame(index = ['Portfolio value', 'Tax liability for full transition', 'Portfolio tracking error'], columns = ['Value']) 
dfPortfolioValueAndTotalTax.loc['Portfolio value', 'Value'] = (dfPortfolioDataFiltered['Number of shares'] * dfPortfolioDataFiltered['Last price']).sum() 
dfPortfolioValueAndTotalTax.loc['Tax liability for full transition', 'Value'] = dfPortfolioDataFiltered['Long-term capital gains'].sum() * longTermTaxRate + dfPortfolioDataFiltered['Short-term capital gains'].sum() * shortTermTaxRate 

startDate = datetime.today() - timedelta(days = noOfYearsToCalculateReplacer * 365) 
endDate = datetime.today() - timedelta(days = 1) 

lstBenchmarkIsins = [] 
lstBenchmarkWeights = [] 
for eachRic, eachWeight in dictBenchmarkData.items(): 
    lstBenchmarkIsins = lstBenchmarkIsins + [dfRicsAndIsins.loc[eachRic, 'ISIN']] 
    lstBenchmarkWeights = lstBenchmarkWeights + [eachWeight] 

dfBenchmarkReturn = benchmark(lstBenchmarkIsins, lstBenchmarkWeights, startDate, endDate) 

dfPortfolioDataReturns = updatePortfolioReturn(list(dfPortfolioDataFiltered['ISIN']) + lstBenchmarkIsins, startDate, endDate) 

# Daily returns for the benchmark 

dfPortfolioBenchmarkAll = dfPortfolioDataReturns.join(dfBenchmarkReturn) 
dfPortfolioBenchmarkReturn = pd.DataFrame(dfPortfolioBenchmarkAll.iloc[:, len(dfPortfolioBenchmarkAll.columns) - 1]) 

currentPortfolioTe = portfolioTrackingError(dfPortfolioDataFiltered.copy(), dfPortfolioDataReturns, dfPortfolioBenchmarkReturn) 

dfPortfolioValueAndTotalTax.loc['Portfolio tracking error', 'Value'] = currentPortfolioTe 

print(dfPortfolioValueAndTotalTax) 

# (ii) Portfolio status for different levels of taxes paid 

dfPortfolioDataFiltered['Position value'] = dfPortfolioDataFiltered['Number of shares'] * dfPortfolioDataFiltered['Last price'] 

dfPortfolioDataFiltered['Tax liability'] = dfPortfolioDataFiltered['Long-term capital gains'] * longTermTaxRate + dfPortfolioDataFiltered['Short-term capital gains'] * shortTermTaxRate 

dfPortfolioDataFiltered['Tax per dollar value'] = dfPortfolioDataFiltered['Tax liability'] / dfPortfolioDataFiltered['Position value'] 

maxTaxPayable = dfPortfolioDataFiltered['Tax liability'].sum() 

totalPortfolioValue = dfPortfolioDataFiltered['Position value'].sum() 

dfTaxLiabilitiesVsProportionTransitionedVsTrackingError = pd.DataFrame(index = range(numOfGroupsForCalculatingTransitionTaxLiabilities + 1), columns = ['Tax liability (%)', 'Tax liability (USD)', 'Value transitioned', 'Proportion transitioned', 'Tracking error']) 
for eachInterval in dfTaxLiabilitiesVsProportionTransitionedVsTrackingError.index: 
    proportionTaxPayment = eachInterval / numOfGroupsForCalculatingTransitionTaxLiabilities 
    
    dfTaxLiabilitiesVsProportionTransitionedVsTrackingError.loc[eachInterval, 'Tax liability (%)'] = proportionTaxPayment 
    
    dfTaxLiabilitiesVsProportionTransitionedVsTrackingError.loc[eachInterval, 'Tax liability (USD)'] = proportionTaxPayment * maxTaxPayable 
    
    dfPortfolioAfterTaxPaymentTemp, dfTransactionsForTaxPaymentTemp = calcPortfolioAfterTaxPayment(dfPortfolioDataFiltered, proportionTaxPayment * maxTaxPayable) 
    
    dfTransactionsForTaxPaymentTemp = dfTransactionsForTaxPaymentTemp[list(dfPortfolioAfterTaxPaymentTemp.columns)] 
    
    dfTaxLiabilitiesVsProportionTransitionedVsTrackingError.loc[eachInterval, 'Value transitioned'] = dfTransactionsForTaxPaymentTemp['Position value'].sum() 
    
    dfTaxLiabilitiesVsProportionTransitionedVsTrackingError.loc[eachInterval, 'Proportion transitioned'] = dfTransactionsForTaxPaymentTemp['Position value'].sum() / totalPortfolioValue 
    
    # Replacing the value transitioned by benchmark constituents in the same proportion as the benchmark 
    for eachIndex in dfClientBenchmark.index: 
        newPositionNumOfShares = round(dfTaxLiabilitiesVsProportionTransitionedVsTrackingError.loc[eachInterval, 'Value transitioned'] * dfClientBenchmark.loc[eachIndex, 'Weights'] / dfClientBenchmark.loc[eachIndex, 'Last price']) 
    
        newPositionLastPrice = dfClientBenchmark.loc[eachIndex, 'Last price'] 
        
        lstToBeBought = [dfClientBenchmark.loc[eachIndex, 'Tickers'], dfClientBenchmark.loc[eachIndex, 'RIC'], dfClientBenchmark.loc[eachIndex, 'ISIN'], datetime.today().strftime("%Y-%m-%d %H:%M:%S"), newPositionNumOfShares, newPositionLastPrice, newPositionLastPrice, dfClientBenchmark.loc[eachIndex, 'GICS sector'], dfClientBenchmark.loc[eachIndex, 'Asset category'], newPositionLastPrice, newPositionLastPrice, 0.0, 0.0, 0.0, newPositionNumOfShares * newPositionLastPrice, 0.0, 0.0] 
        
        serTemp = pd.Series(lstToBeBought, index = dfPortfolioAfterTaxPaymentTemp.columns) 
        
        dfPortfolioAfterTaxPaymentTemp = dfPortfolioAfterTaxPaymentTemp.append(serTemp, ignore_index = True) 
        
        serTemp['Number of shares'] = serTemp['Number of shares'] * -1 
        
        dfTransactionsForTaxPaymentTemp = dfTransactionsForTaxPaymentTemp.append(serTemp, ignore_index = True) 
    
    # Saving portfolio for the selected tax bill proportion for further analysis 
    if proportionTaxPayment == selectedProportion: 
        dfPortfolioSelectedTaxPayment = dfPortfolioAfterTaxPaymentTemp 
        dfTransactionsSelectedTaxPayment = dfTransactionsForTaxPaymentTemp 
    
    # Getting dfPortfolioDataReturnsTemp columns in the same order as the ISINs in dfPortfolioAfterTaxPaymentTemp 
    dfPortfolioDataReturnsTemp = pd.DataFrame() 
    for eachIsin in dfPortfolioAfterTaxPaymentTemp['ISIN']: 
        if dfPortfolioDataReturnsTemp.empty: 
            dfPortfolioDataReturnsTemp = pd.DataFrame(dfPortfolioDataReturns[[eachIsin]].iloc[:, 0]) 
        else: 
            if eachIsin not in list(dfPortfolioDataReturnsTemp.columns): 
                dfPortfolioDataReturnsTemp = dfPortfolioDataReturnsTemp.join(pd.DataFrame(dfPortfolioDataReturns[[eachIsin]].iloc[:, 0])) 
            else: 
                dfTemp = pd.DataFrame(dfPortfolioDataReturns[[eachIsin]].iloc[:, 0]) 
                
                i = 2 
                while eachIsin + ' ' + str(i) in list(dfPortfolioDataReturnsTemp.columns): 
                    i = i + 1 
                    
                dfTemp.columns = [eachIsin + ' ' + str(i)] 

                dfPortfolioDataReturnsTemp = dfPortfolioDataReturnsTemp.join(dfTemp) 
    
    dfPortfolioDataReturnsTemp.columns = list(dfPortfolioAfterTaxPaymentTemp['ISIN']) 
    
    dfTaxLiabilitiesVsProportionTransitionedVsTrackingError.loc[eachInterval, 'Tracking error'] = portfolioTrackingError(dfPortfolioAfterTaxPaymentTemp.copy(), dfPortfolioDataReturnsTemp, dfPortfolioBenchmarkReturn) 
    
print(dfTaxLiabilitiesVsProportionTransitionedVsTrackingError) 

dfOverallTransactionsRequiredForSelectedTaxPayment = dfTransactionsSelectedTaxPayment.groupby(by = ['Security ticker'])[['Number of shares', 'Tax liability']].sum() 

dfOverallTransactionsRequiredForSelectedTaxPayment['Last price'] = np.nan 
dfOverallTransactionsRequiredForSelectedTaxPayment['Turnover'] = np.nan 
dfOverallTransactionsRequiredForSelectedTaxPayment['Trade required'] = '' 
for eachTicker in dfOverallTransactionsRequiredForSelectedTaxPayment.index: 
    dfOverallTransactionsRequiredForSelectedTaxPayment.loc[eachTicker, 'Last price'] =  dfTransactionsSelectedTaxPayment[dfTransactionsSelectedTaxPayment['Security ticker'] == eachTicker]['Last price'].iloc[0] 
    dfOverallTransactionsRequiredForSelectedTaxPayment.loc[eachTicker, 'Turnover'] =  abs(dfOverallTransactionsRequiredForSelectedTaxPayment.loc[eachTicker, 'Last price'] * dfOverallTransactionsRequiredForSelectedTaxPayment.loc[eachTicker, 'Number of shares']) 
    if dfOverallTransactionsRequiredForSelectedTaxPayment.loc[eachTicker, 'Number of shares'] >= 0: 
        dfOverallTransactionsRequiredForSelectedTaxPayment.loc[eachTicker, 'Trade required'] = 'Sell ' + str(dfOverallTransactionsRequiredForSelectedTaxPayment.loc[eachTicker, 'Number of shares']) + ' shares' 
    else: 
        dfOverallTransactionsRequiredForSelectedTaxPayment.loc[eachTicker, 'Trade required'] = 'Buy ' + str(dfOverallTransactionsRequiredForSelectedTaxPayment.loc[eachTicker, 'Number of shares']) + ' shares' 

dfOverallTransactionsRequiredForSelectedTaxPayment = dfOverallTransactionsRequiredForSelectedTaxPayment[['Trade required', 'Last price', 'Tax liability', 'Turnover']] 

print(f'Transactions required for {serTemp} tax payment: ') 
print(dfOverallTransactionsRequiredForSelectedTaxPayment) 

# (iii) Portfolio stats for the selected levels of taxes paid as a proportion of the total 

# Main stats for chosen portfolio 
strTemp = str(int(selectedProportion * 100.0)) + '%' 
indexName = dfTaxLiabilitiesVsProportionTransitionedVsTrackingError[dfTaxLiabilitiesVsProportionTransitionedVsTrackingError['Tax liability (%)'] == selectedProportion].index[0] 
print(f'Tax to be paid as a percentage of the maximum possible tax: {strTemp}') 
print(dfTaxLiabilitiesVsProportionTransitionedVsTrackingError[dfTaxLiabilitiesVsProportionTransitionedVsTrackingError['Tax liability (%)'] == selectedProportion].T[indexName]) 

# Asset class composition 
dfPortfolioDataFilteredTemp = dfPortfolioDataFiltered.copy() 
dfPortfolioDataFilteredTemp['Weight'] = dfPortfolioDataFilteredTemp['Position value'] / dfPortfolioDataFilteredTemp['Position value'].sum() 
dfAssetClassCompositionCurrent = dfPortfolioDataFilteredTemp.groupby(by = ['Asset category'])[['Weight']].sum() 
dfAssetClassCompositionCurrent.columns = ['Current portfolio'] 

dfPortfolioSelectedTaxPayment['Position value'] = dfPortfolioSelectedTaxPayment['Number of shares'] * dfPortfolioSelectedTaxPayment['Last price'] 
dfPortfolioSelectedTaxPayment['Weight'] = dfPortfolioSelectedTaxPayment['Position value'] / dfPortfolioSelectedTaxPayment['Position value'].sum() 
dfAssetClassCompositionSelectedTaxPayment = dfPortfolioSelectedTaxPayment.groupby(by = ['Asset category'])[['Weight']].sum() 
dfAssetClassCompositionSelectedTaxPayment.columns = [f'With {strTemp} tax payment'] 

# SLIDE 2: Calculation of tax liability for each holding and tracking error vs benchmark if that holding is excluded 

lstUniqueSecurities = dfPortfolioDataFiltered['Security ticker'].unique() 
dfConsolidatedPositions = pd.DataFrame(index = lstUniqueSecurities) 

dfConsolidatedPositions = dfConsolidatedPositions.join(pd.DataFrame(dfPortfolioDataFiltered.groupby(by = 'Security ticker')['Number of shares'].sum())) 

dfConsolidatedPositions = dfConsolidatedPositions.join(pd.DataFrame(dfPortfolioDataFiltered.groupby(by = 'Security ticker')['Position value'].sum())) 

dfConsolidatedPositions = dfConsolidatedPositions.join(pd.DataFrame(dfPortfolioDataFiltered.groupby(by = 'Security ticker')['Tax liability'].sum())) 

dfConsolidatedPositions['Portfolio TE on security removal'] = 0.0 
for eachSecurity in dfConsolidatedPositions.index: 
    dfPortfolioDataFilteredTemp = dfPortfolioDataFiltered[dfPortfolioDataFiltered['Security ticker'] != eachSecurity] 
    
    dfPortfolioDataReturnsTemp = dfPortfolioDataReturns[list(dfPortfolioDataFilteredTemp['ISIN'].unique())] 
    
    dfPortfolioBenchmarkAllTemp = dfPortfolioDataReturnsTemp.join(dfBenchmarkReturn) 

    dfPortfolioBenchmarkReturnTemp = pd.DataFrame(dfPortfolioBenchmarkAllTemp.iloc[:, len(dfPortfolioBenchmarkAllTemp.columns) - 1]) 

    securityRemovalTe = portfolioTrackingError(dfPortfolioDataFilteredTemp.copy(), dfPortfolioDataReturnsTemp, dfPortfolioBenchmarkReturnTemp) 
    
    dfConsolidatedPositions.loc[eachSecurity, 'Portfolio TE on security removal'] = securityRemovalTe 

# SLIDE 3: Output format for calcPortfolioLiquidationValue(): dataframe with before tax value, after tax value, tax impact 

dfPortfolioValueOnLiquidation = calcPortfolioLiquidationValue(dfPortfolioDataFiltered.copy(), shortTermTaxRate, longTermTaxRate) 

print(dfPortfolioValueOnLiquidation) 
print('\nPositive value for tax liability means that liquidating the entire portfolio increases its after-tax value') 

# SLIDE 6: Tax liabilities on transitioning different proportions of the portfolio 
# Function withdrawCash() gives two outputs: (1) Portfolio positions after the cash is withdrawn, (2) Transactions required to withdraw the relevant cash amount 

dfPortfolioValuesTransitionedVsTaxLiabilities = pd.DataFrame(index = range(numOfGroupsForCalculatingTransitionTaxLiabilities), columns = ['Proportion transitioned', 'Value transitioned', 'Tax liability (USD)', 'Tracking error']) 
for eachInterval in dfPortfolioValuesTransitionedVsTaxLiabilities.index: 
    proportionTransitioned = (eachInterval + 1) / numOfGroupsForCalculatingTransitionTaxLiabilities 
    
    dfPortfolioValuesTransitionedVsTaxLiabilities.loc[eachInterval, 'Proportion transitioned'] = proportionTransitioned 
    
    dfPortfolioValuesTransitionedVsTaxLiabilities.loc[eachInterval, 'Value transitioned'] = totalPortfolioValue * proportionTransitioned 
    
    dfPortfolioDataFilteredTemp, dfTradesForWithdrawingCashTemp = withdrawCash(dfPortfolioDataFiltered.copy(), dfPortfolioValuesTransitionedVsTaxLiabilities.loc[eachInterval, 'Value transitioned']) 
    
    dfPortfolioValuesTransitionedVsTaxLiabilities.loc[eachInterval, 'Tax liability (USD)'] = dfTradesForWithdrawingCashTemp['Tax liability'].sum() 
    
    # Replacing the value transitioned by benchmark constituents in the same proportion as the benchmark 
    for eachIndex in dfClientBenchmark.index: 
        newPositionNumOfShares = round(dfPortfolioValuesTransitionedVsTaxLiabilities.loc[eachInterval, 'Value transitioned'] * dfClientBenchmark.loc[eachIndex, 'Weights'] / dfClientBenchmark.loc[eachIndex, 'Last price']) 
        
        newPositionLastPrice = dfClientBenchmark.loc[eachIndex, 'Last price'] 
        
        lstToBeBought = [dfClientBenchmark.loc[eachIndex, 'Tickers'], dfClientBenchmark.loc[eachIndex, 'RIC'], dfClientBenchmark.loc[eachIndex, 'ISIN'], datetime.today().strftime("%Y-%m-%d %H:%M:%S"), newPositionNumOfShares, newPositionLastPrice, newPositionLastPrice, dfClientBenchmark.loc[eachIndex, 'GICS sector'], dfClientBenchmark.loc[eachIndex, 'Asset category'], newPositionLastPrice, newPositionLastPrice, 0.0, 0.0, 0.0, newPositionNumOfShares * newPositionLastPrice, 0.0, 0.0] 
        
        serTemp = pd.Series(lstToBeBought, index = dfPortfolioDataFilteredTemp.columns) 
        
        dfPortfolioDataFilteredTemp = dfPortfolioDataFilteredTemp.append(serTemp, ignore_index = True) 
        
        serTemp['Number of shares'] = serTemp['Number of shares'] * -1 
        
        dfTradesForWithdrawingCashTemp = dfTradesForWithdrawingCashTemp.append(serTemp, ignore_index = True) 
    
    # Saving portfolio for the selected assets transition portfolio for further analysis 
    if proportionTransitioned == selectedProportion: 
        dfPortfolioSelectedAssetsTransitioned = dfPortfolioDataFilteredTemp 
        dfTransactionsSelectedAssetsTransitioned = dfTradesForWithdrawingCashTemp 
    
    # Getting dfPortfolioDataReturnsTemp columns in the same order as the ISINs in dfPortfolioDataFilteredTemp 
    dfPortfolioDataReturnsTemp = pd.DataFrame() 
    for eachIsin in dfPortfolioDataFilteredTemp['ISIN']: 
        if dfPortfolioDataReturnsTemp.empty: 
            dfPortfolioDataReturnsTemp = pd.DataFrame(dfPortfolioDataReturns[[eachIsin]].iloc[:, 0]) 
        else: 
            if eachIsin not in list(dfPortfolioDataReturnsTemp.columns): 
                dfPortfolioDataReturnsTemp = dfPortfolioDataReturnsTemp.join(pd.DataFrame(dfPortfolioDataReturns[[eachIsin]].iloc[:, 0])) 
            else: 
                dfTemp = pd.DataFrame(dfPortfolioDataReturns[[eachIsin]].iloc[:, 0]) 
                
                i = 2 
                while eachIsin + ' ' + str(i) in list(dfPortfolioDataReturnsTemp.columns): 
                    i = i + 1 
                    
                dfTemp.columns = [eachIsin + ' ' + str(i)] 

                dfPortfolioDataReturnsTemp = dfPortfolioDataReturnsTemp.join(dfTemp) 
    
    dfPortfolioDataReturnsTemp.columns = list(dfPortfolioDataFilteredTemp['ISIN']) 
    
    dfPortfolioValuesTransitionedVsTaxLiabilities.loc[eachInterval, 'Tracking error'] = portfolioTrackingError(dfPortfolioDataFilteredTemp.copy(), dfPortfolioDataReturnsTemp, dfPortfolioBenchmarkReturn) 
    
print(dfPortfolioValuesTransitionedVsTaxLiabilities) 

# SLIDE 7: Asset class compositions for various portfolios 

dfPortfolioSelectedAssetsTransitioned['Position value'] = dfPortfolioSelectedAssetsTransitioned['Number of shares'] * dfPortfolioSelectedAssetsTransitioned['Last price'] 
dfPortfolioSelectedAssetsTransitioned['Weight'] = dfPortfolioSelectedAssetsTransitioned['Position value'] / dfPortfolioSelectedAssetsTransitioned['Position value'].sum() 

dfAssetClassCompositionSelectedAssetsTransitioned = dfPortfolioSelectedAssetsTransitioned.groupby(by = ['Asset category'])[['Weight']].sum() 
dfAssetClassCompositionSelectedAssetsTransitioned.columns = [f'With {strTemp} assets transitioned'] 

dfAssetClassCompositionBenchmark = dfClientBenchmark.groupby(by = ['Asset category'])[['Weights']].sum() 
dfAssetClassCompositionBenchmark.columns = ['Target portfolio'] 

dfAssetClassCompositionsFinal = dfAssetClassCompositionCurrent.join(dfAssetClassCompositionSelectedTaxPayment).join(dfAssetClassCompositionSelectedAssetsTransitioned).join(dfAssetClassCompositionBenchmark) 
print('Asset class composition: ') 
print(dfAssetClassCompositionsFinal) 

barPlot = (dfAssetClassCompositionsFinal * 100.0).plot.bar() 
barPlot.yaxis.set_major_formatter(mtick.PercentFormatter()) 
plt.xticks(rotation = 0) 

# Save plot  
plt.savefig('Asset_allocation.png') 
plt.show() 

# SLIDE 8: Calculation of losses harvested on the portfolio 

if len(dfPortfolioDataFiltered[dfPortfolioDataFiltered['Tax liability'] < 0]) == 0: 
    lossHarvestingExerciseSuccessful = False 
else: 
    lossHarvestingExerciseSuccessful = True 

    dfPortfolioDataFilteredLossesHarvested = dfPortfolioDataFiltered[dfPortfolioDataFiltered['Tax liability'] >= 0] 
    
    taxSavedLossHarvesting = dfPortfolioDataFiltered[dfPortfolioDataFiltered['Tax liability'] < 0]['Tax liability'].sum() 
    
    dfPortfolioDataReturnsLossesHarvested = dfPortfolioDataReturns[list(dfPortfolioDataFilteredLossesHarvested['ISIN'].unique())] 
        
    dfPortfolioBenchmarkLossesHarvested = dfPortfolioDataReturnsLossesHarvested.join(dfBenchmarkReturn) 
    
    dfPortfolioBenchmarkAllLossesHarvested = pd.DataFrame(dfPortfolioBenchmarkLossesHarvested.iloc[:, len(dfPortfolioBenchmarkLossesHarvested.columns) - 1]) 
    
    lossesHarvestedTe = portfolioTrackingError(dfPortfolioDataFilteredLossesHarvested.copy(), dfPortfolioDataReturnsLossesHarvested, dfPortfolioBenchmarkAllLossesHarvested) 
    
    dfTradesRequiredToHarvestLosses = dfPortfolioDataFiltered[dfPortfolioDataFiltered['Tax liability'] < 0] 
    
    dfTradesRequiredToHarvestLosses = dfTradesRequiredToHarvestLosses.groupby(by = ['Security ticker'])[['Last price', 'Tax liability', 'Number of shares']].sum() 
    
    dfTradesRequiredToHarvestLosses['Trade required'] = '' 
    dfTradesRequiredToHarvestLosses['Turnover'] = 0.0 
    for eachIndex in dfTradesRequiredToHarvestLosses.index: 
        dfTradesRequiredToHarvestLosses.loc[eachIndex, 'Last price'] = dfPortfolioDataFiltered[dfPortfolioDataFiltered['Security ticker'] == eachIndex]['Last price'].iloc[0] 
        dfTradesRequiredToHarvestLosses.loc[eachIndex, 'Trade required'] = 'Sell ' + str(dfTradesRequiredToHarvestLosses.loc[eachIndex, 'Number of shares']) + ' shares' 
        dfTradesRequiredToHarvestLosses.loc[eachIndex, 'Turnover'] = dfTradesRequiredToHarvestLosses.loc[eachIndex, 'Number of shares'] * dfTradesRequiredToHarvestLosses.loc[eachIndex, 'Last price'] 
        
    dfTradesRequiredToHarvestLosses = dfTradesRequiredToHarvestLosses[['Trade required', 'Last price', 'Tax liability', 'Turnover']] 

# SLIDE 9: Impact of selling securities held for the short-term 

dfPortfolioDataFilteredShortTermSecuritiesSold = dfPortfolioDataFiltered[dfPortfolioDataFiltered['Holding period (years)'] <= 1] 

if len(dfPortfolioDataFilteredShortTermSecuritiesSold) == 0: 
    shortTermSecuritiesSellingSuccessful = False 
else: 
    shortTermSecuritiesSellingSuccessful = True 

    dfPortfolioDataFilteredShortTermSecuritiesSold = dfPortfolioDataFiltered[dfPortfolioDataFiltered['Holding period (years)'] <= 1] 
    
    taxImpactShortTermSecuritiesSold = dfPortfolioDataFiltered[dfPortfolioDataFiltered['Short-term capital gains'] != 0]['Tax liability'].sum() 
    
    dfPortfolioDataReturnsShortTermSecuritiesSold = dfPortfolioDataReturns[list(dfPortfolioDataFilteredShortTermSecuritiesSold['ISIN'].unique())] 
        
    dfPortfolioBenchmarkShortTermSecuritiesSold = dfPortfolioDataReturnsShortTermSecuritiesSold.join(dfBenchmarkReturn) 
    
    dfPortfolioBenchmarkAllShortTermSecuritiesSold = pd.DataFrame(dfPortfolioBenchmarkShortTermSecuritiesSold.iloc[:, len(dfPortfolioBenchmarkShortTermSecuritiesSold.columns) - 1]) 
    
    shortTermSecuritiesSoldTe = portfolioTrackingError(dfPortfolioDataFilteredShortTermSecuritiesSold.copy(), dfPortfolioDataReturnsShortTermSecuritiesSold, dfPortfolioBenchmarkAllShortTermSecuritiesSold) 
    
    dfTradesRequiredToSellShortTermSecurities = dfPortfolioDataFiltered[dfPortfolioDataFiltered['Short-term capital gains'] != 0] 
    dfTradesRequiredToSellShortTermSecurities = dfTradesRequiredToSellShortTermSecurities.groupby(by = ['Security ticker'])[['Last price', 'Tax liability', 'Number of shares']].sum() 
    
    dfTradesRequiredToSellShortTermSecurities['Trade required'] = '' 
    dfTradesRequiredToSellShortTermSecurities['Turnover'] = 0.0 
    for eachIndex in dfTradesRequiredToSellShortTermSecurities.index: 
        dfTradesRequiredToSellShortTermSecurities.loc[eachIndex, 'Last price'] = dfPortfolioDataFiltered[dfPortfolioDataFiltered['Security ticker'] == eachIndex]['Last price'].iloc[0] 
        dfTradesRequiredToSellShortTermSecurities.loc[eachIndex, 'Trade required'] = 'Sell ' + str(dfTradesRequiredToSellShortTermSecurities.loc[eachIndex, 'Number of shares']) + ' shares' 
        dfTradesRequiredToSellShortTermSecurities.loc[eachIndex, 'Turnover'] = dfTradesRequiredToSellShortTermSecurities.loc[eachIndex, 'Number of shares'] * dfTradesRequiredToSellShortTermSecurities.loc[eachIndex, 'Last price'] 
        
    dfTradesRequiredToSellShortTermSecurities = dfTradesRequiredToSellShortTermSecurities[['Trade required', 'Last price', 'Tax liability', 'Turnover']] 

# SLIDE 10: Zero tax payment portfolio 
# Here, we sell the portfolio stocks on which there are capital losses, which leads to tax credits 
# We also sell portfolio stocks on which there are capital gains, which offsets the tax credits, resulting in no tax liability 
# While selling stocks with capital gains, stocks on which the tax liability / position value is the lowest are sold first 
# If there are no securities with capital losses, we cannot conduct this exercise 

if len(dfPortfolioDataFiltered[dfPortfolioDataFiltered['Tax liability'] < 0]) == 0 or len(dfPortfolioDataFiltered[dfPortfolioDataFiltered['Tax liability'] > 0]) == 0: 
    lossHarvestingExerciseSuccessful = False 
else: 
    lossHarvestingExerciseSuccessful = True 
    
    dfTradesRequiredForZeroTaxPayment = dfPortfolioDataFiltered[dfPortfolioDataFiltered['Tax liability'] < 0] 
    
    taxSaved = (-1) * dfPortfolioDataFiltered[dfPortfolioDataFiltered['Tax liability'] < 0]['Tax liability'].sum() 

    dfPortfolioDataFilteredZeroTaxPayment = dfPortfolioDataFiltered[dfPortfolioDataFiltered['Tax liability'] >= 0].copy() 
    
    dfPortfolioDataFilteredZeroTaxSortedByTaxPerDollar = dfPortfolioDataFilteredZeroTaxPayment.sort_values(by = 'Tax per dollar value', ascending = True) 
    
    i = 0 
    while taxSaved > 0: 
        relevantIndex = dfPortfolioDataFilteredZeroTaxSortedByTaxPerDollar.index[i] 

        if dfPortfolioDataFilteredZeroTaxSortedByTaxPerDollar.loc[relevantIndex, 'Tax liability'] < taxSaved: 
            dfTradesRequiredForZeroTaxPayment = dfTradesRequiredForZeroTaxPayment.append(dfPortfolioDataFilteredZeroTaxSortedByTaxPerDollar.loc[relevantIndex], ignore_index = True) 
            
            dfPortfolioDataFilteredZeroTaxPayment = dfPortfolioDataFilteredZeroTaxPayment.drop(relevantIndex) 
            
            taxSaved = taxSaved - dfPortfolioDataFilteredZeroTaxSortedByTaxPerDollar.loc[relevantIndex, 'Tax liability'] 
        else: 
            serTemp = dfPortfolioDataFilteredZeroTaxSortedByTaxPerDollar.loc[relevantIndex] 
            
            serTemp['Number of shares'] = round(taxSaved / (serTemp['Last price'] * serTemp['Tax per dollar value'])) 
            
            dfTradesRequiredForZeroTaxPayment = dfTradesRequiredForZeroTaxPayment.append(serTemp, ignore_index = True) 

            dfPortfolioDataFilteredZeroTaxPayment.loc[relevantIndex, 'Number of shares'] = dfPortfolioDataFiltered.loc[relevantIndex, 'Number of shares'] - serTemp['Number of shares'] 
            
            taxSaved = 0 
            
        i = i + 1 
    
    i = 0 
    dfPortfolioDataReturnsZeroTaxPayment = pd.DataFrame() 
    for eachIsin in dfPortfolioDataFilteredZeroTaxPayment['ISIN']: 
        indexToBeIncluded = list(dfPortfolioDataReturns.columns).index(eachIsin) 
        
        if dfPortfolioDataReturnsZeroTaxPayment.empty: 
            dfPortfolioDataReturnsZeroTaxPayment = pd.DataFrame(dfPortfolioDataReturns.iloc[:, indexToBeIncluded]) 
        else: 
            dfPortfolioDataReturnsZeroTaxPayment = dfPortfolioDataReturnsZeroTaxPayment.join(pd.DataFrame(dfPortfolioDataReturns.iloc[:, indexToBeIncluded])) 
            
        lstColumns = list(dfPortfolioDataReturnsZeroTaxPayment.columns) 
        lstColumns[-1] = lstColumns[-1] + ' - ' + str(i) 
        
        dfPortfolioDataReturnsZeroTaxPayment.columns = lstColumns 
        
        i = i + 1 
    
    dfPortfolioDataReturnsZeroTaxPayment.columns = list(dfPortfolioDataFilteredZeroTaxPayment['ISIN']) 
    
    # dfPortfolioDataReturnsZeroTaxPayment = dfPortfolioDataReturns[list(dfPortfolioDataFilteredZeroTaxPayment['ISIN'].unique())] 
        
    dfPortfolioBenchmarkZeroTaxPayment = dfPortfolioDataReturnsZeroTaxPayment.join(dfBenchmarkReturn) 
    
    dfPortfolioBenchmarkAllZeroTaxPayment = pd.DataFrame(dfPortfolioBenchmarkZeroTaxPayment.iloc[:, len(dfPortfolioBenchmarkZeroTaxPayment.columns) - 1]) 
    
    zeroTaxPaymentTe = portfolioTrackingError(dfPortfolioDataFilteredZeroTaxPayment.copy(), dfPortfolioDataReturnsZeroTaxPayment, dfPortfolioBenchmarkAllZeroTaxPayment) 

    dfTradesRequiredForZeroTaxPayment['Tax liability'] = 0.0 
    for eachIndex in dfTradesRequiredForZeroTaxPayment.index: 
        dfTradesRequiredForZeroTaxPayment.loc[eachIndex, 'Tax liability'] = dfTradesRequiredForZeroTaxPayment.loc[eachIndex, 'Number of shares'] * dfTradesRequiredForZeroTaxPayment.loc[eachIndex, 'Last price'] * dfTradesRequiredForZeroTaxPayment.loc[eachIndex, 'Tax per dollar value'] 
    
    dfTradesRequiredForZeroTaxPayment = dfTradesRequiredForZeroTaxPayment.groupby(by = 'Security ticker')['Number of shares', 'Last price', 'Tax liability'].sum() 

    dfTradesRequiredForZeroTaxPayment['Trade required'] = '' 
    dfTradesRequiredForZeroTaxPayment['Last price'] = 0.0 
    dfTradesRequiredForZeroTaxPayment['Turnover'] = 0.0 
    for eachIndex in dfTradesRequiredForZeroTaxPayment.index: 
        dfTradesRequiredForZeroTaxPayment.loc[eachIndex, 'Trade required'] = 'Sell ' + str(dfTradesRequiredForZeroTaxPayment.loc[eachIndex, 'Number of shares']) + ' shares' 
        dfTradesRequiredForZeroTaxPayment.loc[eachIndex, 'Last price'] = dfPortfolioDataFiltered[dfPortfolioDataFiltered['Security ticker'] == eachIndex]['Last price'].iloc[0] 
        dfTradesRequiredForZeroTaxPayment.loc[eachIndex, 'Turnover'] = dfTradesRequiredForZeroTaxPayment.loc[eachIndex, 'Number of shares'] * dfTradesRequiredForZeroTaxPayment.loc[eachIndex, 'Last price'] 

    dfTradesRequiredForZeroTaxPayment = dfTradesRequiredForZeroTaxPayment[['Trade required', 'Last price', 'Tax liability', 'Turnover']] 
    # dfTradesRequiredForZeroTaxPayment = dfTradesRequiredForZeroTaxPayment.set_index('Security ticker') 

# SLIDE 11: Minimum tracking error portfolio 

currentPortfolioTe = portfolioTrackingError(dfPortfolioDataFiltered, dfPortfolioDataReturns, dfPortfolioBenchmarkReturn) 

minPossibleTe, dfWeightsMinTe = calcMinTe(dfPortfolioDataReturns, dfPortfolioBenchmarkReturn, 100, shortTermTaxRate, longTermTaxRate, dfPortfolioDataFiltered) 

dfMinTeOverallWeights = dfWeightsMinTe.groupby(by = 'ISIN').sum() 

dfMinTrackingErrorPortfolio, dfTaxImpactForMinTePortfolio = calcMinTaxPaymentForTargetWeights(dfPortfolioDataFiltered, dfMinTeOverallWeights) 

dfMinTrackingErrorPortfolioGrouped = dfMinTrackingErrorPortfolio.groupby(by = 'Security ticker')[['Number of shares']].sum() 

dfCurrentPortfolioGrouped = dfPortfolioDataFiltered.groupby(by = 'Security ticker')[['Number of shares']].sum() 

dfTradesRequiredTemp = dfMinTrackingErrorPortfolioGrouped - dfCurrentPortfolioGrouped 

dfTradesRequiredForMinTePortfolio = pd.DataFrame(index = dfTradesRequiredTemp.index, columns = ['Trade required', 'Last price', 'Tax liability', 'Turnover']) 
for eachTicker in dfTradesRequiredForMinTePortfolio.index: 
    if dfTradesRequiredTemp.loc[eachTicker, 'Number of shares'] < 0: 
        dfTradesRequiredForMinTePortfolio.loc[eachTicker, 'Trade required'] = 'Sell ' + str((-1) * dfTradesRequiredTemp.loc[eachTicker, 'Number of shares']) + ' shares' 
    else: 
        dfTradesRequiredForMinTePortfolio.loc[eachTicker, 'Trade required'] = 'Buy ' + str(dfTradesRequiredTemp.loc[eachTicker, 'Number of shares']) + ' shares' 
        
    dfTradesRequiredForMinTePortfolio.loc[eachTicker, 'Last price'] = dfPortfolioDataFiltered[dfPortfolioDataFiltered['Security ticker'] == eachTicker]['Last price'].iloc[0] 

    dfTradesRequiredForMinTePortfolio.loc[eachTicker, 'Tax liability'] = dfTaxImpactForMinTePortfolio.loc[eachTicker, 'USD value'] 

    dfTradesRequiredForMinTePortfolio.loc[eachTicker, 'Turnover'] = abs(dfTradesRequiredTemp.loc[eachTicker, 'Number of shares']) * dfTradesRequiredForMinTePortfolio.loc[eachTicker, 'Last price'] 

# SLIDE 12: Replacement ETFs: Only spyders corresponding to the relevant GICS sectors are chosen for replacement; methodology for fixed income ETFs yet to be decided 
# Create a dictionary sectorDict that matches each GICS sector to the relevant ETF ticker 

sectorToEtfDict =   { 
                        "Information Technology": "XLK", 
                        "Industrials": "XLI", 
                        "Financials": "XLF", 
                        "Communication Services": "XLC", 
                        "Real Estate": "XLRE", 
                        "Energy": "XLE", 
                        "Consumer Discretionary": "XLY", 
                        "Materials": "XLB", 
                        "Health Care": "XLV", 
                        "Utilities": "XLU", 
                        "Consumer Staples": "XLP" 
                    } 

dfPortfolioComponentsReplacementTickers = dfPortfolioDataFiltered.groupby(by = 'Security ticker')['Security ticker', 'GICS sector'].first() 
dfPortfolioComponentsReplacementTickersTemp = dfPortfolioDataFiltered.groupby(by = 'Security ticker')['Number of shares', 'Position value', 'Tax liability'].sum() 

dfPortfolioComponentsReplacementTickers['Replacement security'] = '' 
for eachTicker in dfPortfolioComponentsReplacementTickers.index: 
    if dfPortfolioComponentsReplacementTickers.loc[eachTicker, 'GICS sector'] in list(sectorToEtfDict.keys()): 
        dfPortfolioComponentsReplacementTickers.loc[eachTicker, 'Replacement security'] = sectorToEtfDict[dfPortfolioComponentsReplacementTickers.loc[eachTicker, 'GICS sector']] 
    else: 
        dfPortfolioComponentsReplacementTickers.loc[eachTicker, 'Replacement security'] = 'AGG' 

dfPortfolioComponentsReplacementTickers = dfPortfolioComponentsReplacementTickers.join(dfPortfolioComponentsReplacementTickersTemp) 

# dfReplacementSecurityDetails = ek.get_data(list(dfPortfolioComponentsReplacementTickers['Replacement security'].values), fields = ['DSPLY_NAME', 'TR.PriceCloseDate', 'TR.PriceClose'])[0] 
dfReplacementSecurityDetails = pd.read_excel('20220206 Replacement security details.xlsx', header = 0, index_col = 0) 

dfReplacementSecurityDetails = dfReplacementSecurityDetails.groupby(by = 'Instrument')[dfReplacementSecurityDetails.columns[1:]].first() 

dfReplacementSecurityDetails = dfReplacementSecurityDetails.rename(columns = {'DSPLY_NAME' : 'Name'}) 

dfPortfolioComponentsReplacementTickers['Replacement security price'] = 0.0 
dfPortfolioComponentsReplacementTickers['Replacement security name'] = '' 
dfPortfolioComponentsReplacementTickers['Trades required'] = '' 
for eachTicker in dfPortfolioComponentsReplacementTickers.index: 
    dfPortfolioComponentsReplacementTickers.loc[eachTicker, 'Replacement security price'] = dfReplacementSecurityDetails.loc[dfPortfolioComponentsReplacementTickers.loc[eachTicker, 'Replacement security'], 'Price Close'] 
    dfPortfolioComponentsReplacementTickers.loc[eachTicker, 'Replacement security name'] = dfReplacementSecurityDetails.loc[dfPortfolioComponentsReplacementTickers.loc[eachTicker, 'Replacement security'], 'Name'] 
    dfPortfolioComponentsReplacementTickers.loc[eachTicker, 'Trades required'] = 'Sell ' + str(dfPortfolioComponentsReplacementTickers.loc[eachTicker, 'Number of shares']) + ' ' + eachTicker + ', buy ' + str(round(dfPortfolioComponentsReplacementTickers.loc[eachTicker, 'Position value'] / dfPortfolioComponentsReplacementTickers.loc[eachTicker, 'Replacement security price'])) + ' ' + dfPortfolioComponentsReplacementTickers.loc[eachTicker, 'Replacement security'] 

dfPortfolioComponentsReplacements = dfPortfolioComponentsReplacementTickers[['GICS sector', 'Replacement security', 'Trades required', 'Tax liability']] 

# SLIDE 13: Most tax-efficient way of withdrawing money 
# Here, we are assuming that the transactions to settle the tax liability would be done later, i.e. when the individual files taxes 
# Function withdrawCash() gives two outputs: (1) Portfolio positions after the cash is withdrawn, (2) Transactions required to withdraw the relevant cash amount 

dfPortfolioDataFilteredWithdrawCash, dfTradesForWithdrawingCash = withdrawCash(dfPortfolioDataFiltered.copy(), amountToBeWithdrawn) 

i = 0 
dfPortfolioDataReturnsWithdrawCash = pd.DataFrame() 
for eachIsin in dfPortfolioDataFilteredWithdrawCash['ISIN']: 
    indexToBeIncluded = list(dfPortfolioDataReturns.columns).index(eachIsin) 
        
    if dfPortfolioDataReturnsWithdrawCash.empty: 
        dfPortfolioDataReturnsWithdrawCash = pd.DataFrame(dfPortfolioDataReturns.iloc[:, indexToBeIncluded]) 
    else: 
        dfPortfolioDataReturnsWithdrawCash = dfPortfolioDataReturnsWithdrawCash.join(pd.DataFrame(dfPortfolioDataReturns.iloc[:, indexToBeIncluded])) 
        
    lstColumns = list(dfPortfolioDataReturnsWithdrawCash.columns) 
    lstColumns[-1] = lstColumns[-1] + ' - ' + str(i) 
        
    dfPortfolioDataReturnsWithdrawCash.columns = lstColumns 
        
    i = i + 1 
    
dfPortfolioDataReturnsWithdrawCash.columns = list(dfPortfolioDataFilteredWithdrawCash['ISIN']) 

dfPortfolioBenchmarkWithdrawCash = dfPortfolioDataReturnsWithdrawCash.join(dfBenchmarkReturn) 

dfPortfolioBenchmarkAllWithdrawCash = pd.DataFrame(dfPortfolioBenchmarkWithdrawCash.iloc[:, len(dfPortfolioBenchmarkWithdrawCash.columns) - 1]) 

withdrawCashTe = portfolioTrackingError(dfPortfolioDataFilteredWithdrawCash.copy(), dfPortfolioDataReturnsWithdrawCash, dfPortfolioBenchmarkAllWithdrawCash) 

dfTradesForWithdrawingCash = dfTradesForWithdrawingCash.groupby(by = 'Security ticker')['Number of shares', 'Tax liability'].sum() 

dfTradesForWithdrawingCash['Trade required'] = '' 
dfTradesForWithdrawingCash['Last price'] = 0.0 
# dfTradesForWithdrawingCash['Tax liability'] = 0.0 
dfTradesForWithdrawingCash['Turnover'] = 0.0 
for eachTicker in dfTradesForWithdrawingCash.index: 
    dfTradesForWithdrawingCash.loc[eachTicker, 'Trade required'] = 'Sell ' + str(dfTradesForWithdrawingCash.loc[eachTicker, 'Number of shares']) + ' shares' 
    dfTradesForWithdrawingCash.loc[eachTicker, 'Last price'] = dfPortfolioDataFiltered[dfPortfolioDataFiltered['Security ticker'] == eachTicker]['Last price'].iloc[0] 
    # dfTradesForWithdrawingCash.loc[eachTicker, 'Tax liability'] = dfTradesForWithdrawingCash.loc[eachTicker, 'Last price'] * dfTradesForWithdrawingCash.loc[eachTicker, 'Number of shares'] * dfTradesForWithdrawingCash.loc[eachTicker, 'Tax per dollar value'] 
    dfTradesForWithdrawingCash.loc[eachTicker, 'Turnover'] = dfTradesForWithdrawingCash.loc[eachTicker, 'Number of shares'] * dfTradesForWithdrawingCash.loc[eachTicker, 'Last price'] 

dfTradesForWithdrawingCash = dfTradesForWithdrawingCash[['Trade required', 'Last price', 'Tax liability', 'Turnover']] 

# SLIDE 14: Replacement of a portfolio security by another security 

componentReplacementTaxLiability = dfPortfolioDataFiltered[dfPortfolioDataFiltered['ISIN'] == tickerToBeReplaced]['Tax liability'].sum() 

dfPortfolioDataFilteredComponentReplaced = dfPortfolioDataFiltered[dfPortfolioDataFiltered['ISIN'] != tickerToBeReplaced] 

serReplacementEntry = pd.Series(index = dfPortfolioDataFilteredComponentReplaced.columns) 

# dfReplacementTickerData = ek.get_data([replacementTicker], fields = ['TR.ISIN','TR.PriceClose.date', 'TR.PriceCloseDate', 'TR.PriceClose', 'TR.GICSSector'])[0] 
dfReplacementTickerData = pd.read_excel('20220206 Replacement ticker data.xlsx', header = 0, index_col = 0) 

# dfReplacementTickerTimeSeries = ek.get_data([replacementTicker], fields = ['TR.PriceClose.date', 'TR.PriceClose'], parameters = {'Frq': 'D', 'SDate': startDate.strftime('%Y-%m-%d'), 'CALCMETHOD': 'CLOSE', 'EDate': endDate.strftime('%Y-%m-%d')})[0] 
dfReplacementTickerTimeSeries = pd.read_excel('20220206 Replacement ticker time series.xlsx', header = 0, index_col = 0) 

dfReplacementTickerTimeSeries = dfReplacementTickerTimeSeries.drop(columns = ['Instrument']) 

dfReplacementTickerTimeSeries = dfReplacementTickerTimeSeries.set_index('Date') 

dfReplacementTickerTimeSeries.index = pd.to_datetime(dfReplacementTickerTimeSeries.index).date 

# dfReplacementTickerTimeSeries = dfReplacementTickerTimeSeries.loc[dfPortfolioDataReturns.index[]] 

serReplacementEntry['Security ticker'] = replacementTicker 
serReplacementEntry['Purchase date'] = datetime.today().date() 
serReplacementEntry['Last price'] = dfReplacementTickerData.loc[dfReplacementTickerData.index[0], 'Price Close'] 
serReplacementEntry['Number of shares'] = round(dfPortfolioDataFiltered[dfPortfolioDataFiltered['Security ticker'] == tickerToBeReplaced]['Position value'].sum() / serReplacementEntry['Last price']) 
serReplacementEntry['ISIN'] = dfReplacementTickerData.loc[dfReplacementTickerData.index[0], 'ISIN'] 
serReplacementEntry['GICS sector'] = dfReplacementTickerData.loc[dfReplacementTickerData.index[0], 'GICS Sector Name'] 
serReplacementEntry['Purchase date prices'] = dfReplacementTickerData.loc[dfReplacementTickerData.index[0], 'Price Close'] 
serReplacementEntry['Final purchase price'] = dfReplacementTickerData.loc[dfReplacementTickerData.index[0], 'Price Close'] 
serReplacementEntry['Holding period (years)'] = 0.0 
serReplacementEntry['Long-term capital gains'] = 0.0 
serReplacementEntry['Short-term capital gains'] = 0.0 
serReplacementEntry['Tax liability'] = 0.0 
serReplacementEntry['Tax per dollar value'] = 0.0 
serReplacementEntry['Position value'] = serReplacementEntry['Last price'] * serReplacementEntry['Number of shares'] 
serReplacementEntry['Current weight'] = 0.0 

dfReplacementTickerTimeSeries.columns = [serReplacementEntry['ISIN']] 

dfPortfolioDataFilteredComponentReplaced = dfPortfolioDataFilteredComponentReplaced.append(serReplacementEntry, ignore_index = True) 

i = 0 
dfPortfolioDataReturnsComponentReplaced = pd.DataFrame() 
for eachIsin in dfPortfolioDataFilteredComponentReplaced['ISIN']: 
    if eachIsin in list(dfPortfolioDataReturns.columns): 
        indexToBeIncluded = list(dfPortfolioDataReturns.columns).index(eachIsin) 
        
        if dfPortfolioDataReturnsComponentReplaced.empty: 
            dfPortfolioDataReturnsComponentReplaced = pd.DataFrame(dfPortfolioDataReturns.iloc[:, indexToBeIncluded]) 
        else: 
            dfPortfolioDataReturnsComponentReplaced = dfPortfolioDataReturnsComponentReplaced.join(pd.DataFrame(dfPortfolioDataReturns.iloc[:, indexToBeIncluded])) 
            
        lstColumns = list(dfPortfolioDataReturnsComponentReplaced.columns) 
        lstColumns[-1] = lstColumns[-1] + ' - ' + str(i) 
            
        dfPortfolioDataReturnsComponentReplaced.columns = lstColumns 
    else: 
        if dfPortfolioDataReturnsWithdrawCash.empty: 
            dfPortfolioDataReturnsComponentReplaced = dfReplacementTickerTimeSeries 
        else: 
            dfPortfolioDataReturnsComponentReplaced = dfPortfolioDataReturnsComponentReplaced.join(dfReplacementTickerTimeSeries) 
        
    i = i + 1 
    
dfPortfolioDataReturnsComponentReplaced.columns = list(dfPortfolioDataFilteredComponentReplaced['ISIN'].values) 

dfPortfolioBenchmarkComponentReplaced = dfPortfolioDataReturnsComponentReplaced.join(dfBenchmarkReturn) 

dfPortfolioBenchmarkAllComponentReplaced = pd.DataFrame(dfPortfolioBenchmarkComponentReplaced.iloc[:, len(dfPortfolioBenchmarkComponentReplaced.columns) - 1]) 

componentReplacedTe = portfolioTrackingError(dfPortfolioDataFilteredComponentReplaced.copy(), dfPortfolioDataReturnsComponentReplaced, dfPortfolioBenchmarkAllComponentReplaced) 

dfTradesForComponentReplacement = pd.DataFrame(index = [tickerToBeReplaced, replacementTicker], columns = ['Trade required', 'Last price', 'Tax liability']) 

dfTradesForComponentReplacement['Trade required'] = '' 
dfTradesForComponentReplacement['Last price'] = 0.0 
# dfTradesForComponentReplacement['Tax liability'] = 0.0 
dfTradesForComponentReplacement['Turnover'] = 0.0 
for eachTicker in dfTradesForComponentReplacement.index: 
    dfTradesForComponentReplacement.loc[eachTicker, 'Last price'] = dfPortfolioDataFilteredComponentReplaced[dfPortfolioDataFilteredComponentReplaced['Security ticker'] == eachTicker]['Last price'].iloc[0] 

    if eachTicker == tickerToBeReplaced: 
        dfTradesForComponentReplacement.loc[eachTicker, 'Trade required'] = 'Sell ' + str(dfPortfolioDataFiltered[dfPortfolioDataFiltered['Security ticker'] == eachTicker]['Number of shares'].sum()) + ' shares' 
        dfTradesForComponentReplacement.loc[eachTicker, 'Tax liability'] = dfPortfolioDataFiltered[dfPortfolioDataFiltered['Security ticker'] == eachTicker]['Tax liability'].sum() 
        dfTradesForComponentReplacement.loc[eachTicker, 'Turnover'] = dfPortfolioDataFiltered[dfPortfolioDataFiltered['Security ticker'] == eachTicker]['Number of shares'].sum() * dfTradesForComponentReplacement.loc[eachTicker, 'Last price'] 
    else: 
        dfTradesForComponentReplacement.loc[eachTicker, 'Trade required'] = 'Buy ' + str(serReplacementEntry['Number of shares']) + ' shares' 
        dfTradesForComponentReplacement.loc[eachTicker, 'Tax liability'] = 0.0 
        dfTradesForComponentReplacement.loc[eachTicker, 'Turnover'] = serReplacementEntry['Number of shares'] * dfTradesForComponentReplacement.loc[eachTicker, 'Last price'] 

dfTradesForWithdrawingCash = dfTradesForWithdrawingCash[['Trade required', 'Last price', 'Tax liability', 'Turnover']] 


# SLIDE 15: Calculating the tax frontier 

# numOfOptimizationPoints = 10 
# lstOptimizationPoints = list(np.linspace(minPossibleTe, currentPortfolioTe, numOfOptimizationPoints)) 


#--------------------------------------------------------------
# OUTPUT USING STREAMLIT 
#--------------------------------------------------------------

st.title('Tax transition analysis for ' + clientName) 

st.header('Portfolio constituents and tracking error analysis: ') 

dfTemp = pd.DataFrame(currentPortfolioTe, index = ['Current portfolio'], columns = ['Tracking error (%)']) 
dfTemp = changeDataframeFormat(dfTemp, ['{:.2%}']) 

st.dataframe(dfTemp) 

dfTemp = changeDataframeFormat(dfConsolidatedPositions.copy(), ['{:,.0f}', '${:,.0f}', '${:,.0f}', '{:.2%}']) 

st.dataframe(dfTemp) 

st.header('Impact of portfolio liquidation: ') 

dfTemp = pd.DataFrame(dfPortfolioValueOnLiquidation.loc['Tax liability', 'USD value'] / dfPortfolioValueOnLiquidation.loc['Before tax value', 'USD value'], index = ['Tax liability on selling entire portfolio'], columns = ['Percentage of portfolio value']) 
dfTemp = changeDataframeFormat(dfTemp, ['{:.2%}']) 

st.dataframe(dfTemp) 

dfTemp = pd.DataFrame(0, index = ['Short-term gains', 'Short-term losses', 'Long-term gains', 'Long-term losses', 'Total'], columns = ['Unrealized P/L', 'Tax liability']) 
dfTemp.loc['Short-term gains', 'Unrealized P/L'] = dfPortfolioDataFiltered['Short-term capital gains'][dfPortfolioDataFiltered['Short-term capital gains'] >= 0].sum() 
dfTemp.loc['Short-term gains', 'Tax liability'] = shortTermTaxRate * dfTemp.loc['Short-term gains', 'Unrealized P/L'] 
dfTemp.loc['Short-term losses', 'Unrealized P/L'] = dfPortfolioDataFiltered['Short-term capital gains'][dfPortfolioDataFiltered['Short-term capital gains'] < 0].sum() 
dfTemp.loc['Short-term losses', 'Tax liability'] = shortTermTaxRate * dfTemp.loc['Short-term losses', 'Unrealized P/L'] 
dfTemp.loc['Long-term gains', 'Unrealized P/L'] = dfPortfolioDataFiltered['Long-term capital gains'][dfPortfolioDataFiltered['Short-term capital gains'] >= 0].sum() 
dfTemp.loc['Long-term gains', 'Tax liability'] = longTermTaxRate * dfTemp.loc['Long-term gains', 'Unrealized P/L'] 
dfTemp.loc['Long-term losses', 'Unrealized P/L'] = dfPortfolioDataFiltered['Long-term capital gains'][dfPortfolioDataFiltered['Short-term capital gains'] < 0].sum() 
dfTemp.loc['Long-term losses', 'Tax liability'] = longTermTaxRate * dfTemp.loc['Long-term losses', 'Unrealized P/L'] 
dfTemp.loc['Total', 'Unrealized P/L'] = dfPortfolioDataFiltered['Long-term capital gains'].sum() + dfPortfolioDataFiltered['Short-term capital gains'].sum() 
dfTemp.loc['Total', 'Tax liability'] = longTermTaxRate * dfPortfolioDataFiltered['Long-term capital gains'].sum() + shortTermTaxRate * dfPortfolioDataFiltered['Short-term capital gains'].sum() 

dfTemp = changeDataframeFormat(dfTemp.copy(), ['${:,.0f}', '${:,.0f}']) 

st.dataframe(dfTemp) 

st.header('Current portfolio statistics versus alternative portfolios: ') 

dfTemp = dfPortfolioValueAndTotalTax.copy() 
# dfTemp.loc['Portfolio tracking error', 'Value'] = dfTemp.loc['Portfolio tracking error', 'Value'] 
dfTemp.columns = ['Current portfolio'] 

dfTemp = changeDataframeFormat(dfTemp.T.copy(), ['${:,.0f}', '${:,.0f}', '{:.2%}']) 

st.dataframe(dfTemp.T) 

dfTemp = pd.DataFrame(dfTaxLiabilitiesVsProportionTransitionedVsTrackingError[dfTaxLiabilitiesVsProportionTransitionedVsTrackingError['Tax liability (%)'] == selectedProportion].T[indexName]).copy() 
strColName = f'Target portfolio ({selectedProportion * 100}% tax bill)' 
dfTemp.columns = [strColName] 

dfTemp = changeDataframeFormat(dfTemp.T.copy(), ['{:.0%}', '${:,.0f}', '${:,.0f}', '{:.2%}', '{:.2%}']) 

st.dataframe(dfTemp.T) 

dfTemp = dfPortfolioValuesTransitionedVsTaxLiabilities[dfPortfolioValuesTransitionedVsTaxLiabilities['Proportion transitioned'] == selectedProportion].T.copy() 
strColName = f'Target portfolio ({selectedProportion * 100}% value transitioned)' 
dfTemp.columns = [strColName] 

dfTemp = changeDataframeFormat(dfTemp.T.copy(), ['{:.0%}', '${:,.0f}', '${:,.0f}', '{:.2%}']) 

st.dataframe(dfTemp.T) 

st.header('Portfolio proportions transitioned for different tax bills: ') 

dfTemp = changeDataframeFormat(dfTaxLiabilitiesVsProportionTransitionedVsTrackingError.copy(), ['{:.0%}', '${:,.0f}', '${:,.0f}', '{:.2%}', '{:.2%}']) 

st.dataframe(dfTemp) 

st.header('Tax liabilities for transitioning different portfolio proportions: ') 

dfTemp = changeDataframeFormat(dfPortfolioValuesTransitionedVsTaxLiabilities.copy(), ['{:.0%}', '${:,.0f}', '${:,.0f}', '{:.2%}']) 

st.dataframe(dfTemp) 

st.header('Asset allocations for different variations: ') 

dfTemp = changeDataframeFormat(dfAssetClassCompositionsFinal.copy(), ['{:.2%}', '{:.2%}', '{:.2%}', '{:.2%}']) 

st.dataframe(dfTemp) 

# st.bar_chart(dfAssetClassCompositionsFinal) 

# barPlot = (dfAssetClassCompositionsFinal * 100.0).plot(kind = 'bar') 
# st.pyplot(barPlot) 

imageDisplay = Image.open('Asset_allocation.png') 

st.image(imageDisplay) 

st.header('Harvesting losses on the portfolio: ') 

if lossHarvestingExerciseSuccessful == False: 
    st.write('NO SECURITIES WITH LOSSES IN PORTFOLIO') 
else: 
    dfTemp = pd.DataFrame(0, index = ['Tax liability (USD)', 'Tracking error', 'Total turnover'], columns = ['Value']) 
    
    dfTemp.loc['Tax liability (USD)', 'Value'] = taxSavedLossHarvesting 
    dfTemp.loc['Tracking error', 'Value'] = lossesHarvestedTe 
    dfTemp.loc['Total turnover', 'Value'] = dfTradesRequiredToHarvestLosses['Turnover'].sum() 
    
    dfTemp = changeDataframeFormat(dfTemp.T.copy(), ['${:,.0f}', '{:.2%}', '${:,.0f}']) 

    st.dataframe(dfTemp.T) 
    
    if st.checkbox('Show trades required to harvest losses'): 
        st.write('Trades required to harvest losses: ') 
        
        dfTemp = changeDataframeFormat(dfTradesRequiredToHarvestLosses.copy(), ['', '{:,.2f}', '${:,.0f}', '${:,.0f}']) 

        st.dataframe(dfTemp) 

st.header('Selling securities held for <= 1 year: ') 

if shortTermSecuritiesSellingSuccessful == False: 
    st.write('NO SECURITIES IN PORTFOLIO WITH HOLDING PERIOD <= 1 YEAR') 
else: 
    dfTemp = pd.DataFrame(0, index = ['Tax liability (USD)', 'Tracking error', 'Total turnover'], columns = ['Value']) 
    
    dfTemp.loc['Tax liability (USD)', 'Value'] = taxImpactShortTermSecuritiesSold 
    dfTemp.loc['Tracking error', 'Value'] = shortTermSecuritiesSoldTe 
    dfTemp.loc['Total turnover', 'Value'] = dfTradesRequiredToSellShortTermSecurities['Turnover'].sum() 
        
    dfTemp = changeDataframeFormat(dfTemp.T.copy(), ['${:,.0f}', '{:.2%}', '${:,.0f}']) 

    st.dataframe(dfTemp.T) 
    
    if st.checkbox('Show trades required to sell securities held for <= 1 year'): 
        st.write('Trades required to sell securities held for <= 1 year: ') 
        # st.dataframe(dfTradesRequiredToSellShortTermSecurities.style.format('{:.2f}')) 

        dfTemp = changeDataframeFormat(dfTradesRequiredToSellShortTermSecurities.copy(), ['', '{:,.2f}', '${:,.0f}', '${:,.0f}']) 

        st.dataframe(dfTemp) 

st.header('Portfolio transactions for 0 net tax payment: ') 

if lossHarvestingExerciseSuccessful == False: 
    st.write('NO SECURITIES WITH LOSSES IN PORTFOLIO') 
else: 
    dfTemp = pd.DataFrame(0, index = ['Tax liability (USD)', 'Tracking error', 'Total turnover'], columns = ['Value']) 
    
    dfTemp.loc['Tax liability (USD)', 'Value'] = 0 
    dfTemp.loc['Tracking error', 'Value'] = zeroTaxPaymentTe 
    dfTemp.loc['Total turnover', 'Value'] = dfTradesRequiredToSellShortTermSecurities['Turnover'].sum() 
    
    dfTemp = changeDataframeFormat(dfTemp.T.copy(), ['${:,.0f}', '{:.2%}', '${:,.0f}']) 

    st.dataframe(dfTemp.T) 
    
    if st.checkbox('Show trades required for 0 net tax payment'): 
        st.write('Trades required for 0 net tax payment: ') 

        dfTemp = changeDataframeFormat(dfTradesRequiredForZeroTaxPayment.copy(), ['', '{:,.2f}', '${:,.0f}', '${:,.0f}']) 

        st.dataframe(dfTemp) 

st.header('Portfolio transactions for achieving minimum tracking error: ') 

dfTemp = pd.DataFrame(0, index = ['Tax liability (USD)', 'Tracking error', 'Total turnover'], columns = ['Value']) 
    
dfTemp.loc['Tax liability (USD)', 'Value'] = dfTaxImpactForMinTePortfolio['USD value'].sum() 
dfTemp.loc['Tracking error', 'Value'] = minPossibleTe 
dfTemp.loc['Total turnover', 'Value'] = dfTradesRequiredForMinTePortfolio['Turnover'].sum() 

dfTemp = changeDataframeFormat(dfTemp.T.copy(), ['${:,.0f}', '{:.2%}', '${:,.0f}']) 

st.dataframe(dfTemp.T) 
    
if st.checkbox('Show trades required for minimum tracking error'): 
    st.write('Trades required for minimum tracking error: ') 

    dfTemp = changeDataframeFormat(dfTradesRequiredForMinTePortfolio.copy(), ['', '{:,.2f}', '${:,.0f}', '${:,.0f}']) 

    st.dataframe(dfTemp) 

st.header('Replacement of portfolio securities by GICS sector ETFs: ') 

dfTemp = changeDataframeFormat(dfPortfolioComponentsReplacements.copy(), ['', '', '', '${:,.0f}']) 

st.dataframe(dfTemp) 

st.header('Tax-efficient way to withdraw ' + "${:,.0f}".format(amountToBeWithdrawn) + ': ') 

dfTemp = pd.DataFrame(0, index = ['Tax liability (USD)', 'Tracking error', 'Total turnover'], columns = ['Value']) 
    
dfTemp.loc['Tax liability (USD)', 'Value'] = dfTradesForWithdrawingCash['Tax liability'].sum() 
dfTemp.loc['Tracking error', 'Value'] = withdrawCashTe 
dfTemp.loc['Total turnover', 'Value'] = dfTradesForWithdrawingCash['Turnover'].sum() 

dfTemp = changeDataframeFormat(dfTemp.T.copy(), ['${:,.0f}', '{:.2%}', '${:,.0f}']) 

st.dataframe(dfTemp.T) 
    
if st.checkbox('Show trades required to withdraw ' + "${:,.0f}".format(amountToBeWithdrawn) + ' in a cash efficient manner'): 
    st.write('Trades required to withdraw ' + "${:,.0f}".format(amountToBeWithdrawn) + ' in a cash efficient manner') 

    dfTemp = changeDataframeFormat(dfTradesForWithdrawingCash.copy(), ['', '{:,.2f}', '${:,.0f}', '${:,.0f}']) 

    st.dataframe(dfTemp) 

st.header('Replacement of ' + tickerToBeReplaced + ' by ' + replacementTicker + ': ') 

dfTemp = pd.DataFrame(0, index = ['Tax liability (USD)', 'Tracking error', 'Total turnover'], columns = ['Value']) 
    
dfTemp.loc['Tax liability (USD)', 'Value'] = dfTradesForComponentReplacement['Tax liability'].sum() 
dfTemp.loc['Tracking error', 'Value'] = withdrawCashTe 
dfTemp.loc['Total turnover', 'Value'] = dfTradesForComponentReplacement['Turnover'].sum() 

dfTemp = changeDataframeFormat(dfTemp.T.copy(), ['${:,.0f}', '{:.2%}', '${:,.0f}']) 

st.dataframe(dfTemp.T) 
    
if st.checkbox('Show trades required to replace ' + tickerToBeReplaced + ' with ' + replacementTicker): 
    st.write('Trades required to replace ' + tickerToBeReplaced + ' with ' + replacementTicker + ': ') 

    dfTemp = changeDataframeFormat(dfTradesForComponentReplacement.copy(), ['', '{:,.2f}', '${:,.0f}', '${:,.0f}']) 

    st.dataframe(dfTradesForComponentReplacement) 
    
boolCreatePresentation = st.button('Create presentation') 

#--------------------------------------------------------------
# GENERATING PPTS 
#--------------------------------------------------------------

if boolCreatePresentation == True: 
    # Creating blank slide with Yayati logo 
    
    yayatiLogo = 'logo.png' 
    
    prs = Presentation() 
    
    blankSlideLayout = prs.slide_layouts[5] 
    
    # Slide 1 
    
    slide = prs.slides.add_slide(prs.slide_layouts[0]) 
    shapes = slide.shapes 
    title = slide.shapes.title 
    title.text = "Tax strategy" 
    
    bodyShape = shapes.placeholders[0]
    tf = bodyShape.text_frame 
    tf.text = 'Tax transition analysis for ' + clientName 
    
    titlePara = slide.shapes.title.text_frame.paragraphs[0]
    titlePara.font.name = 'Calibri'
    titlePara.font.size = Pt(25)
    
    left = Inches(0.2) 
    top = Inches(0.2) 
    height = Inches(0.5) 
    pic = slide.shapes.add_picture(yayatiLogo, left, top, height = height) 
    
    # Slide 2 
    
    slide = prs.slides.add_slide(prs.slide_layouts[5]) 
    title = slide.shapes.title 
    
    title.text = 'Portfolio constituents and tracking error analysis' 
    
    titlePara = slide.shapes.title.text_frame.paragraphs[0] 
    titlePara.font.name = 'Calibri' 
    titlePara.font.size = Pt(25) 
    
    left = Inches(0.2) 
    top = Inches(0.2) 
    height = Inches(0.5) 
    pic = slide.shapes.add_picture(yayatiLogo, left, top, height = height) 
    
    shape = slide.shapes.add_table(2, 2, Inches(2), Inches(1.5), Inches(6), Inches(1)) 
    table = shape.table 
    table.cell(0, 1).text = 'Tracking error (%)' 
    table.cell(0, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    
    table.cell(1, 0).text = 'Current portfolio' 
    
    table.cell(1, 1).text = str(round(currentPortfolioTe * 100.0, 2)) + '%' 
    table.cell(1, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    
    numOfPositions = len(dfConsolidatedPositions.index) 
    
    shape = slide.shapes.add_table(numOfPositions + 1, 5, Inches(0.5), Inches(3), Inches(9), Inches(numOfPositions / 2)) 
    table = shape.table 
    
    i = 1 
    for eachCol in dfConsolidatedPositions.columns: 
        table.cell(0, i).text = eachCol 
        table.cell(0, i).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
        
        i = i + 1 
    
    i = 1 
    for eachIndex in dfConsolidatedPositions.index: 
        table.cell(i, 0).text = eachIndex 
    
        j = 1 
        for eachCol in dfConsolidatedPositions.columns: 
            if j == 2 or j == 3: 
                table.cell(i, j).text = "${:,.0f}".format(dfConsolidatedPositions.loc[eachIndex, eachCol]) 
            elif j == 4: 
                table.cell(i, j).text = str(round(dfConsolidatedPositions.loc[eachIndex, eachCol] * 100.0, 2)) + '%' 
            else: 
                table.cell(i, j).text = "{:,.0f}".format(dfConsolidatedPositions.loc[eachIndex, eachCol]) 
            
            table.cell(i, j).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
            
            j = j + 1 
            
        i = i + 1 
    
    # Slide 3 
    
    slide = prs.slides.add_slide(prs.slide_layouts[5]) 
    title = slide.shapes.title 
    
    title.text = 'Impact of portfolio liquidation' 
    
    titlePara = slide.shapes.title.text_frame.paragraphs[0] 
    titlePara.font.name = 'Calibri' 
    titlePara.font.size = Pt(25) 
    
    left = Inches(0.2) 
    top = Inches(0.2) 
    height = Inches(0.5) 
    pic = slide.shapes.add_picture(yayatiLogo, left, top, height = height) 
    
    shape = slide.shapes.add_table(6, 3, Inches(1), Inches(1.5), Inches(8), Inches(3)) 
    table = shape.table 
    table.cell(0, 1).text = 'Unrealized P/L' 
    table.cell(0, 2).text = 'Tax liability' 
    table.cell(0, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    table.cell(0, 2).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    
    table.cell(1, 0).text = 'Short-term gains' 
    table.cell(2, 0).text = 'Short-term losses' 
    table.cell(3, 0).text = 'Long-term gains' 
    table.cell(4, 0).text = 'Long-term losses' 
    table.cell(5, 0).text = 'Total' 
    
    table.cell(1, 1).text = "${:,.0f}".format(dfPortfolioDataFiltered['Short-term capital gains'][dfPortfolioDataFiltered['Short-term capital gains'] > 0].sum()) 
    table.cell(1, 2).text = "${:,.0f}".format(shortTermTaxRate * dfPortfolioDataFiltered['Short-term capital gains'][dfPortfolioDataFiltered['Short-term capital gains'] > 0].sum()) 
    table.cell(1, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    table.cell(1, 2).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    
    table.cell(2, 1).text = "${:,.0f}".format(dfPortfolioDataFiltered['Short-term capital gains'][dfPortfolioDataFiltered['Short-term capital gains'] < 0].sum()) 
    table.cell(2, 2).text = "${:,.0f}".format(shortTermTaxRate * dfPortfolioDataFiltered['Short-term capital gains'][dfPortfolioDataFiltered['Short-term capital gains'] < 0].sum()) 
    table.cell(2, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    table.cell(2, 2).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    
    table.cell(3, 1).text = "${:,.0f}".format(dfPortfolioDataFiltered['Long-term capital gains'][dfPortfolioDataFiltered['Long-term capital gains'] > 0].sum()) 
    table.cell(3, 2).text = "${:,.0f}".format(longTermTaxRate * dfPortfolioDataFiltered['Long-term capital gains'][dfPortfolioDataFiltered['Long-term capital gains'] > 0].sum()) 
    table.cell(3, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    table.cell(3, 2).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    
    table.cell(4, 1).text = "${:,.0f}".format(dfPortfolioDataFiltered['Long-term capital gains'][dfPortfolioDataFiltered['Long-term capital gains'] < 0].sum()) 
    table.cell(4, 2).text = "${:,.0f}".format(longTermTaxRate * dfPortfolioDataFiltered['Long-term capital gains'][dfPortfolioDataFiltered['Long-term capital gains'] < 0].sum()) 
    table.cell(4, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    table.cell(4, 2).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    
    table.cell(5, 1).text = "${:,.0f}".format(dfPortfolioDataFiltered['Long-term capital gains'].sum() + dfPortfolioDataFiltered['Short-term capital gains'].sum()) 
    table.cell(5, 2).text = "${:,.0f}".format(longTermTaxRate * dfPortfolioDataFiltered['Long-term capital gains'].sum() + shortTermTaxRate * dfPortfolioDataFiltered['Short-term capital gains'].sum()) 
    table.cell(5, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    table.cell(5, 2).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    
    shape = slide.shapes.add_table(2, 2, Inches(2), Inches(5), Inches(6), Inches(1.5)) 
    table = shape.table 
    table.cell(0, 1).text = 'Percentage of portfolio value' 
    table.cell(0, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    
    table.cell(1, 0).text = 'Tax liability on selling entire portfolio' 
    
    table.cell(1, 1).text = str(round(dfPortfolioValueOnLiquidation.loc['Tax liability', 'USD value'] / dfPortfolioValueOnLiquidation.loc['Before tax value', 'USD value'] * 100.0, 2)) + '%' 
    table.cell(1, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    
    # Slide 4 
    
    slide = prs.slides.add_slide(prs.slide_layouts[5]) 
    title = slide.shapes.title 
    
    title.text = 'Current portfolio statistics versus alternative portfolios' 
    
    titlePara = slide.shapes.title.text_frame.paragraphs[0] 
    titlePara.font.name = 'Calibri' 
    titlePara.font.size = Pt(25) 
    
    left = Inches(0.2) 
    top = Inches(0.2) 
    height = Inches(0.5) 
    pic = slide.shapes.add_picture(yayatiLogo, left, top, height = height) 
    
    shape = slide.shapes.add_table(4, 2, Inches(0.5), Inches(1.5), Inches(9), Inches(1.6)) 
    table = shape.table 
    table.cell(0, 1).text = 'Current portfolio' 
    table.cell(0, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    
    table.cell(1, 0).text = dfPortfolioValueAndTotalTax.index[0] 
    table.cell(2, 0).text = dfPortfolioValueAndTotalTax.index[1] 
    table.cell(3, 0).text = dfPortfolioValueAndTotalTax.index[2] 
    
    table.cell(1, 1).text = "${:,.0f}".format(dfPortfolioValueAndTotalTax.loc[dfPortfolioValueAndTotalTax.index[0], 'Value']) 
    table.cell(1, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    
    table.cell(2, 1).text = "${:,.0f}".format(dfPortfolioValueAndTotalTax.loc[dfPortfolioValueAndTotalTax.index[1], 'Value']) 
    table.cell(2, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    
    table.cell(3, 1).text = "{:,.2f}".format(dfPortfolioValueAndTotalTax.loc[dfPortfolioValueAndTotalTax.index[2], 'Value'] * 100.0) + '%' 
    table.cell(3, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    
    serTransitionedPortfolioTaxBill = dfTaxLiabilitiesVsProportionTransitionedVsTrackingError[dfTaxLiabilitiesVsProportionTransitionedVsTrackingError['Tax liability (%)'] == selectedProportion].T[indexName] 
    
    shape = slide.shapes.add_table(5, 2, Inches(0.5), Inches(3.3), Inches(9), Inches(1.5)) 
    table = shape.table 
    table.cell(0, 1).text = f'Target portfolio ({strTemp} tax bill)' 
    table.cell(0, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    
    i = 0 
    for eachIndex in serTransitionedPortfolioTaxBill.iloc[1:].index: 
        table.cell(i + 1, 0).text = serTransitionedPortfolioTaxBill.index[i + 1] 
        table.cell(i + 1, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT 
        
        if eachIndex == 'Tax liability (USD)' or eachIndex == 'Value transitioned': 
            table.cell(i + 1, 1).text = "${:,.0f}".format(serTransitionedPortfolioTaxBill[eachIndex]) 
        else: 
            table.cell(i + 1, 1).text = "{:,.2f}".format(serTransitionedPortfolioTaxBill[eachIndex] * 100.0) + '%' 
        
        table.cell(i + 1, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
        
        i = i + 1 
    
    dfTransitionedPortfolioValueTransitioned = dfPortfolioValuesTransitionedVsTaxLiabilities[dfPortfolioValuesTransitionedVsTaxLiabilities['Proportion transitioned'] == selectedProportion].T 
    colName = dfTransitionedPortfolioValueTransitioned.columns[0] 
    serTransitionedPortfolioValueTransitioned = dfTransitionedPortfolioValueTransitioned[colName] 
    
    shape = slide.shapes.add_table(4, 2, Inches(0.5), Inches(5.5), Inches(9), Inches(1.5)) 
    table = shape.table 
    table.cell(0, 1).text = f'Target portfolio ({strTemp} value transitioned)' 
    table.cell(0, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    
    i = 0 
    for eachIndex in serTransitionedPortfolioValueTransitioned.index[1:]: 
        table.cell(i + 1, 0).text = serTransitionedPortfolioValueTransitioned.index[i] 
        table.cell(i + 1, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT 
        
        if eachIndex == 'Tax liability (USD)' or eachIndex == 'Value transitioned': 
            table.cell(i + 1, 1).text = "${:,.0f}".format(serTransitionedPortfolioValueTransitioned[eachIndex]) 
        else: 
            table.cell(i + 1, 1).text = "{:,.2f}".format(serTransitionedPortfolioValueTransitioned[eachIndex] * 100.0) + '%' 
        
        table.cell(i + 1, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
        
        i = i + 1 
    
    # Slide 5 
    
    slide = prs.slides.add_slide(prs.slide_layouts[5]) 
    title = slide.shapes.title 
    
    title.text = 'Portfolio proportions transitioned for different tax bills' 
    
    titlePara = slide.shapes.title.text_frame.paragraphs[0] 
    titlePara.font.name = 'Calibri' 
    titlePara.font.size = Pt(25) 
    
    left = Inches(0.2) 
    top = Inches(0.2) 
    height = Inches(0.5) 
    pic = slide.shapes.add_picture(yayatiLogo, left, top, height = height) 
    
    shape = slide.shapes.add_table(numOfGroupsForCalculatingTransitionTaxLiabilities + 2, 6, Inches(0.5), Inches(1.5), Inches(9), Inches(numOfGroupsForCalculatingTransitionTaxLiabilities / 2 + 0.5)) 
    table = shape.table 
    i = 0 
    for eachColumn in dfTaxLiabilitiesVsProportionTransitionedVsTrackingError.columns: 
        table.cell(0, i + 1).text = eachColumn 
        
        i = i + 1 
    
    i = 0 
    for eachIndex in dfTaxLiabilitiesVsProportionTransitionedVsTrackingError.index: 
        table.cell(i + 1, 0).text = str(i + 1) 
    
        table.cell(i + 1, 1).text = "{:,.0f}".format(dfTaxLiabilitiesVsProportionTransitionedVsTrackingError.loc[eachIndex, 'Tax liability (%)'] * 100.0) + '%' 
        table.cell(i + 1, 2).text = "${:,.0f}".format(dfTaxLiabilitiesVsProportionTransitionedVsTrackingError.loc[eachIndex, 'Tax liability (USD)']) 
        table.cell(i + 1, 3).text = "${:,.0f}".format(dfTaxLiabilitiesVsProportionTransitionedVsTrackingError.loc[eachIndex, 'Value transitioned']) 
        table.cell(i + 1, 4).text = "{:,.2f}".format(dfTaxLiabilitiesVsProportionTransitionedVsTrackingError.loc[eachIndex, 'Proportion transitioned'] * 100.0) + '%' 
        table.cell(i + 1, 5).text = "{:,.2f}".format(dfTaxLiabilitiesVsProportionTransitionedVsTrackingError.loc[eachIndex, 'Tracking error'] * 100.0) + '%' 
    
        table.cell(i + 1, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
        table.cell(i + 1, 2).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
        table.cell(i + 1, 3).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
        table.cell(i + 1, 4).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
        table.cell(i + 1, 5).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
        
        i = i + 1 
    
    # Slide 6 
    
    slide = prs.slides.add_slide(prs.slide_layouts[5]) 
    title = slide.shapes.title 
    
    title.text = 'Tax liabilities for transitioning different portfolio proportions' 
    
    titlePara = slide.shapes.title.text_frame.paragraphs[0] 
    titlePara.font.name = 'Calibri' 
    titlePara.font.size = Pt(25) 
    
    left = Inches(0.2) 
    top = Inches(0.2) 
    height = Inches(0.5) 
    pic = slide.shapes.add_picture(yayatiLogo, left, top, height = height) 
    
    shape = slide.shapes.add_table(numOfGroupsForCalculatingTransitionTaxLiabilities + 1, 5, Inches(1), Inches(1.5), Inches(8), Inches(numOfGroupsForCalculatingTransitionTaxLiabilities / 2)) 
    table = shape.table 
    table.cell(0, 1).text = 'Proportion transitioned' 
    table.cell(0, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    table.cell(0, 2).text = 'Value transitioned' 
    table.cell(0, 2).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    table.cell(0, 3).text = 'Tax liability (USD)' 
    table.cell(0, 3).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    table.cell(0, 4).text = 'Tracking error' 
    table.cell(0, 4).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    
    i = 0 
    for eachIndex in dfPortfolioValuesTransitionedVsTaxLiabilities.index: 
        table.cell(i + 1, 0).text = str(i + 1) 
    
        table.cell(i + 1, 1).text = "{:,.0f}".format(dfPortfolioValuesTransitionedVsTaxLiabilities.loc[eachIndex, 'Proportion transitioned'] * 100.0) + '%' 
        table.cell(i + 1, 2).text = "${:,.0f}".format(dfPortfolioValuesTransitionedVsTaxLiabilities.loc[eachIndex, 'Value transitioned']) 
        table.cell(i + 1, 3).text = "${:,.0f}".format(dfPortfolioValuesTransitionedVsTaxLiabilities.loc[eachIndex, 'Tax liability (USD)']) 
        if dfPortfolioValuesTransitionedVsTaxLiabilities.loc[eachIndex, 'Tracking error'] == 'NA': 
            table.cell(i + 1, 4).text = dfPortfolioValuesTransitionedVsTaxLiabilities.loc[eachIndex, 'Tracking error'] 
        else: 
            table.cell(i + 1, 4).text = "{:,.2f}".format(dfPortfolioValuesTransitionedVsTaxLiabilities.loc[eachIndex, 'Tracking error'] * 100.0) + '%' 
    
        table.cell(i + 1, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
        table.cell(i + 1, 2).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
        table.cell(i + 1, 3).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
        table.cell(i + 1, 4).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
        
        i = i + 1 
    
    # Slide 7 
     
    slide = prs.slides.add_slide(prs.slide_layouts[5]) 
    title = slide.shapes.title 
    
    title.text = 'Asset allocations for different variations' 
    
    titlePara = slide.shapes.title.text_frame.paragraphs[0] 
    titlePara.font.name = 'Calibri' 
    titlePara.font.size = Pt(25) 
    
    left = Inches(0.2) 
    top = Inches(0.2) 
    height = Inches(0.5) 
    pic = slide.shapes.add_picture(yayatiLogo, left, top, height = height) 
        
    shape = slide.shapes.add_table(3, 5, Inches(0.5), Inches(1.5), Inches(9), Inches(len(dfAssetClassCompositionsFinal.index) / 2)) 
    table = shape.table 
    i = 0 
    for eachColumn in dfAssetClassCompositionsFinal.columns: 
        table.cell(0, i + 1).text = eachColumn 
        table.cell(0, i + 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
        
        i = i + 1 
    
    i = 0 
    for eachIndex in dfAssetClassCompositionsFinal.index: 
        table.cell(i + 1, 0).text = eachIndex 
        table.cell(i + 1, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT 
        
        j = 0 
        for eachColumn in dfAssetClassCompositionsFinal.columns: 
            table.cell(i + 1, j + 1).text = "{:,.2f}".format(dfAssetClassCompositionsFinal.loc[eachIndex, eachColumn] * 100.0) + '%' 
            table.cell(i + 1, j + 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
            
            j = j + 1 
    
        i = i + 1 
    
    pic = slide.shapes.add_picture('Asset_allocation.png', Inches(2), Inches(3.3), Inches(6)) 
    
    # Slide 8 
    
    slide = prs.slides.add_slide(prs.slide_layouts[5]) 
    title = slide.shapes.title 
    
    title.text = 'Harvesting losses on the portfolio' 
    
    titlePara = slide.shapes.title.text_frame.paragraphs[0] 
    titlePara.font.name = 'Calibri' 
    titlePara.font.size = Pt(25) 
    
    left = Inches(0.2) 
    top = Inches(0.2) 
    height = Inches(0.5) 
    pic = slide.shapes.add_picture(yayatiLogo, left, top, height = height) 
    
    if shortTermSecuritiesSellingSuccessful == False: 
        shape = slide.shapes.add_table(1, 1, Inches(2), Inches(1.5), Inches(6), Inches(1.6)) 
        table = shape.table 
        table.cell(0, 0).text = 'NO SECURITIES IN PORTFOLIO WITH HOLDING PERIOD <= 1 YEAR' 
        table.cell(0, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER 
    else:  
        shape = slide.shapes.add_table(4, 2, Inches(2), Inches(1.5), Inches(6), Inches(1.6)) 
        table = shape.table 
        table.cell(0, 1).text = 'Value' 
        table.cell(0, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
        table.cell(1, 0).text = 'Tax liability (USD)' 
        table.cell(2, 0).text = 'Tracking error' 
        table.cell(3, 0).text = 'Total turnover' 
        
        table.cell(1, 1).text = "${:,.0f}".format(taxSavedLossHarvesting) 
        table.cell(1, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
        table.cell(2, 1).text = str(round(lossesHarvestedTe * 100.0, 2)) + '%' 
        table.cell(2, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
        table.cell(3, 1).text = "${:,.0f}".format(dfTradesRequiredToHarvestLosses['Turnover'].sum()) 
        table.cell(3, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
        
        # dfTradesRequiredToHarvestLosses[['Security ticker', 'Trade required', 'Price', 'Loss harvested']] 
        numOfTradesRequiredToHarvestLosses = len(dfTradesRequiredToHarvestLosses.index) 
        
        shape = slide.shapes.add_table(numOfTradesRequiredToHarvestLosses + 1, 5, Inches(0.5), Inches(3.5), Inches(9), Inches(numOfTradesRequiredToHarvestLosses / 2)) 
        table = shape.table 
        
        i = 1 
        for eachCol in dfTradesRequiredToHarvestLosses.columns: 
            table.cell(0, i).text = eachCol 
            table.cell(0, i).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
            
            i = i + 1 
        
        i = 1 
        for eachIndex in dfTradesRequiredToHarvestLosses.index: 
            table.cell(i, 0).text = eachIndex 
        
            j = 1 
            for eachCol in dfTradesRequiredToHarvestLosses.columns: 
                if j == 1: 
                    table.cell(i, j).text = dfTradesRequiredToHarvestLosses.loc[eachIndex, eachCol] 
                elif j == 2: 
                    table.cell(i, j).text = str(round(dfTradesRequiredToHarvestLosses.loc[eachIndex, eachCol], 2)) 
                else: 
                    table.cell(i, j).text = "${:,.0f}".format(dfTradesRequiredToHarvestLosses.loc[eachIndex, eachCol]) 
                
                table.cell(i, j).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
        
                j = j + 1 
                
            i = i + 1 
    
    # Slide 9 
    
    slide = prs.slides.add_slide(prs.slide_layouts[5]) 
    title = slide.shapes.title 
    
    title.text = 'Selling securities held for <= 1 year' 
    
    titlePara = slide.shapes.title.text_frame.paragraphs[0] 
    titlePara.font.name = 'Calibri' 
    titlePara.font.size = Pt(25) 
    
    left = Inches(0.2) 
    top = Inches(0.2) 
    height = Inches(0.5) 
    pic = slide.shapes.add_picture(yayatiLogo, left, top, height = height) 
    
    if shortTermSecuritiesSellingSuccessful == False: 
        shape = slide.shapes.add_table(1, 1, Inches(2), Inches(1.5), Inches(6), Inches(1.6)) 
        table = shape.table 
        table.cell(0, 0).text = 'NO SECURITIES HELD FOR LESS THAN OR EQUAL TO 1 YEAR' 
        table.cell(0, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER 
    else:  
        shape = slide.shapes.add_table(4, 2, Inches(2), Inches(1.5), Inches(6), Inches(1.6)) 
        table = shape.table 
        table.cell(0, 1).text = 'Value' 
        table.cell(0, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
        table.cell(1, 0).text = 'Tax liability (USD)' 
        table.cell(2, 0).text = 'Tracking error' 
        table.cell(3, 0).text = 'Total turnover' 
        
        table.cell(1, 1).text = "${:,.0f}".format(taxImpactShortTermSecuritiesSold) 
        table.cell(1, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
        table.cell(2, 1).text = str(round(shortTermSecuritiesSoldTe * 100.0, 2)) + '%' 
        table.cell(2, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
        table.cell(3, 1).text = "${:,.0f}".format(dfTradesRequiredToSellShortTermSecurities['Turnover'].sum()) 
        table.cell(3, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
        
        # dfTradesRequiredToHarvestLosses[['Security ticker', 'Trade required', 'Price', 'Loss harvested']] 
        numOfTradesRequiredToSellShortTermSecurities = len(dfTradesRequiredToSellShortTermSecurities.index) 
        
        shape = slide.shapes.add_table(numOfTradesRequiredToSellShortTermSecurities + 1, 5, Inches(0.5), Inches(3.5), Inches(9), Inches(numOfTradesRequiredToSellShortTermSecurities / 2)) 
        table = shape.table 
        
        i = 1 
        for eachCol in dfTradesRequiredToSellShortTermSecurities.columns: 
            table.cell(0, i).text = eachCol 
            table.cell(0, i).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
            
            i = i + 1 
        
        i = 1 
        for eachIndex in dfTradesRequiredToSellShortTermSecurities.index: 
            table.cell(i, 0).text = eachIndex 
        
            j = 1 
            for eachCol in dfTradesRequiredToSellShortTermSecurities.columns: 
                if j == 1: 
                    table.cell(i, j).text = dfTradesRequiredToSellShortTermSecurities.loc[eachIndex, eachCol] 
                elif j == 2: 
                    table.cell(i, j).text = str(round(dfTradesRequiredToSellShortTermSecurities.loc[eachIndex, eachCol], 2)) 
                else: 
                    table.cell(i, j).text = "${:,.0f}".format(dfTradesRequiredToSellShortTermSecurities.loc[eachIndex, eachCol]) 
                
                table.cell(i, j).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
        
                j = j + 1 
                
            i = i + 1 
    
    # Slide 10 
    
    slide = prs.slides.add_slide(prs.slide_layouts[5]) 
    title = slide.shapes.title 
    
    title.text = 'Portfolio transactions for 0 net tax payment' 
    
    titlePara = slide.shapes.title.text_frame.paragraphs[0] 
    titlePara.font.name = 'Calibri' 
    titlePara.font.size = Pt(25) 
    
    left = Inches(0.2) 
    top = Inches(0.2) 
    height = Inches(0.5) 
    pic = slide.shapes.add_picture(yayatiLogo, left, top, height = height) 
    
    if lossHarvestingExerciseSuccessful == False: 
        shape = slide.shapes.add_table(1, 1, Inches(2), Inches(1.5), Inches(6), Inches(1.6)) 
        table = shape.table 
        table.cell(0, 0).text = 'NOT POSSIBLE TO MAKE TRANSACTIONS WITHOUT ANY TAX IMPACT' 
        table.cell(0, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER 
    else:  
        shape = slide.shapes.add_table(4, 2, Inches(2), Inches(1.5), Inches(6), Inches(1.6)) 
        table = shape.table 
        table.cell(0, 1).text = 'Value' 
        table.cell(0, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
        table.cell(1, 0).text = 'Tax liability (USD)' 
        table.cell(2, 0).text = 'Tracking error' 
        table.cell(3, 0).text = 'Total turnover' 
        
        table.cell(1, 1).text = "${:,.0f}".format(0) 
        table.cell(1, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
        table.cell(2, 1).text = str(round(zeroTaxPaymentTe * 100.0, 2)) + '%' 
        table.cell(2, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
        table.cell(3, 1).text = "${:,.0f}".format(dfTradesRequiredForZeroTaxPayment['Turnover'].sum()) 
        table.cell(3, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
        
        # dfTradesRequiredToHarvestLosses[['Security ticker', 'Trade required', 'Price', 'Loss harvested']] 
        numOfTradesRequiredForZeroTaxPayment = len(dfTradesRequiredForZeroTaxPayment.index) 
        
        shape = slide.shapes.add_table(numOfTradesRequiredForZeroTaxPayment + 1, 5, Inches(0.5), Inches(3.5), Inches(9), Inches(numOfTradesRequiredForZeroTaxPayment / 2)) 
        table = shape.table 
        
        i = 1 
        for eachCol in dfTradesRequiredForZeroTaxPayment.columns: 
            table.cell(0, i).text = eachCol 
            table.cell(0, i).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
            
            i = i + 1 
        
        i = 1 
        for eachIndex in dfTradesRequiredForZeroTaxPayment.index: 
            table.cell(i, 0).text = eachIndex 
        
            j = 1 
            for eachCol in dfTradesRequiredForZeroTaxPayment.columns: 
                if j == 1: 
                    table.cell(i, j).text = dfTradesRequiredForZeroTaxPayment.loc[eachIndex, eachCol] 
                elif j == 2: 
                    table.cell(i, j).text = str(round(dfTradesRequiredForZeroTaxPayment.loc[eachIndex, eachCol], 2)) 
                else: 
                    table.cell(i, j).text = "${:,.0f}".format(dfTradesRequiredForZeroTaxPayment.loc[eachIndex, eachCol]) 
                
                table.cell(i, j).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
        
                j = j + 1 
                
            i = i + 1 
    
    # Slide 11 
    
    slide = prs.slides.add_slide(prs.slide_layouts[5]) 
    title = slide.shapes.title 
    
    title.text = 'Portfolio transactions for achieving minimum tracking error' 
    
    titlePara = slide.shapes.title.text_frame.paragraphs[0] 
    titlePara.font.name = 'Calibri' 
    titlePara.font.size = Pt(25) 
    
    left = Inches(0.2) 
    top = Inches(0.2) 
    height = Inches(0.5) 
    pic = slide.shapes.add_picture(yayatiLogo, left, top, height = height) 
    
    shape = slide.shapes.add_table(4, 2, Inches(2), Inches(1.5), Inches(6), Inches(1.6)) 
    table = shape.table 
    table.cell(0, 1).text = 'Value' 
    table.cell(0, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    table.cell(1, 0).text = 'Tax liability (USD)' 
    table.cell(2, 0).text = 'Tracking error' 
    table.cell(3, 0).text = 'Total turnover' 
    
    table.cell(1, 1).text = "${:,.0f}".format(dfTaxImpactForMinTePortfolio['USD value'].sum()) 
    table.cell(1, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    table.cell(2, 1).text = str(round(minPossibleTe * 100.0, 2)) + '%' 
    table.cell(2, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    table.cell(3, 1).text = "${:,.0f}".format(dfTradesRequiredForMinTePortfolio['Turnover'].sum()) 
    table.cell(3, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
        
    numOfTradesRequiredForMinTePortfolio = len(dfTradesRequiredForMinTePortfolio.index) 
        
    shape = slide.shapes.add_table(numOfTradesRequiredForMinTePortfolio + 1, 5, Inches(0.5), Inches(3.5), Inches(9), Inches(numOfTradesRequiredForMinTePortfolio / 2)) 
    table = shape.table 
    
    i = 1 
    for eachCol in dfTradesRequiredForMinTePortfolio.columns: 
        table.cell(0, i).text = eachCol 
        table.cell(0, i).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
           
        i = i + 1 
        
    i = 1 
    for eachIndex in dfTradesRequiredForMinTePortfolio.index: 
        table.cell(i, 0).text = eachIndex 
        
        j = 1 
        for eachCol in dfTradesRequiredForMinTePortfolio.columns: 
            if j == 1: 
                table.cell(i, j).text = dfTradesRequiredForMinTePortfolio.loc[eachIndex, eachCol] 
            elif j == 2: 
                table.cell(i, j).text = str(round(dfTradesRequiredForMinTePortfolio.loc[eachIndex, eachCol], 2)) 
            else: 
                table.cell(i, j).text = "${:,.0f}".format(dfTradesRequiredForMinTePortfolio.loc[eachIndex, eachCol]) 
                
            table.cell(i, j).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
        
            j = j + 1 
                
        i = i + 1 
    
    # Slide 12 
    
    slide = prs.slides.add_slide(prs.slide_layouts[5]) 
    title = slide.shapes.title 
    
    title.text = 'Replacement of portfolio securities by GICS sector ETFs' 
    
    titlePara = slide.shapes.title.text_frame.paragraphs[0] 
    titlePara.font.name = 'Calibri' 
    titlePara.font.size = Pt(25) 
    
    left = Inches(0.2) 
    top = Inches(0.2) 
    height = Inches(0.5) 
    pic = slide.shapes.add_picture(yayatiLogo, left, top, height = height) 
    
    numOfRows = len(dfPortfolioComponentsReplacements.index) 
    numOfCols = len(dfPortfolioComponentsReplacements.columns) 
    
    shape = slide.shapes.add_table(numOfRows + 1, numOfCols + 1, Inches(0.5), Inches(1.3), Inches(9), Inches(1.5)) 
    table = shape.table 
    i = 1 
    for eachCol in dfPortfolioComponentsReplacements.columns: 
        table.cell(0, i).text = eachCol 
        
        i = i + 1 
    
    i = 1 
    for eachRow in dfPortfolioComponentsReplacements.index: 
        table.cell(i, 0).text = eachRow 
        
        j = 1 
        for eachCol in dfPortfolioComponentsReplacements.columns: 
            if j < 4: 
                table.cell(i, j).text = '' if dfPortfolioComponentsReplacements.loc[eachRow, eachCol] == None else dfPortfolioComponentsReplacements.loc[eachRow, eachCol] 
            else: 
                table.cell(i, j).text = "${:,.0f}".format(dfPortfolioComponentsReplacements.loc[eachRow, eachCol]) 
                
            j = j + 1 
        
        i = i + 1 
    
    # Slide 13 
    
    slide = prs.slides.add_slide(prs.slide_layouts[5]) 
    title = slide.shapes.title 
    
    title.text = 'Tax-efficient way to withdraw ' + "${:,.0f}".format(amountToBeWithdrawn) 
    
    titlePara = slide.shapes.title.text_frame.paragraphs[0] 
    titlePara.font.name = 'Calibri' 
    titlePara.font.size = Pt(25) 
    
    left = Inches(0.2) 
    top = Inches(0.2) 
    height = Inches(0.5) 
    pic = slide.shapes.add_picture(yayatiLogo, left, top, height = height) 
    
    shape = slide.shapes.add_table(4, 2, Inches(2), Inches(1.5), Inches(6), Inches(1.6)) 
    table = shape.table 
    table.cell(0, 1).text = 'Value' 
    table.cell(0, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    table.cell(1, 0).text = 'Tax liability (USD)' 
    table.cell(2, 0).text = 'Tracking error' 
    table.cell(3, 0).text = 'Total turnover' 
    
    table.cell(1, 1).text = "${:,.0f}".format(dfTradesForWithdrawingCash['Tax liability'].sum()) 
    table.cell(1, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    table.cell(2, 1).text = str(round(withdrawCashTe * 100.0, 2)) + '%' 
    table.cell(2, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    table.cell(3, 1).text = "${:,.0f}".format(dfTradesForWithdrawingCash['Turnover'].sum()) 
    table.cell(3, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    
    numOfTradesRequiredForWithdrawingCash = len(dfTradesForWithdrawingCash.index) 
    
    shape = slide.shapes.add_table(numOfTradesRequiredForWithdrawingCash + 1, 5, Inches(0.5), Inches(3.5), Inches(9), Inches(numOfTradesRequiredForWithdrawingCash / 2)) 
    table = shape.table 
    
    i = 1 
    for eachCol in dfTradesForWithdrawingCash.columns: 
        table.cell(0, i).text = eachCol 
        table.cell(0, i).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
        
        i = i + 1 
        
    i = 1 
    for eachIndex in dfTradesForWithdrawingCash.index: 
        table.cell(i, 0).text = eachIndex 
        
        j = 1 
        for eachCol in dfTradesForWithdrawingCash.columns: 
            if j == 1: 
                table.cell(i, j).text = dfTradesForWithdrawingCash.loc[eachIndex, eachCol] 
            elif j == 2: 
                table.cell(i, j).text = str(round(dfTradesForWithdrawingCash.loc[eachIndex, eachCol], 2)) 
            else: 
                table.cell(i, j).text = "${:,.0f}".format(dfTradesForWithdrawingCash.loc[eachIndex, eachCol]) 
                
            table.cell(i, j).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
        
            j = j + 1 
                
        i = i + 1 
    
    # Slide 14 
    
    slide = prs.slides.add_slide(prs.slide_layouts[5]) 
    title = slide.shapes.title 
    
    title.text = 'Replacement of ' + tickerToBeReplaced + ' by ' + replacementTicker 
    
    titlePara = slide.shapes.title.text_frame.paragraphs[0] 
    titlePara.font.name = 'Calibri' 
    titlePara.font.size = Pt(25) 
    
    left = Inches(0.2) 
    top = Inches(0.2) 
    height = Inches(0.5) 
    pic = slide.shapes.add_picture(yayatiLogo, left, top, height = height) 
    
    shape = slide.shapes.add_table(4, 2, Inches(2), Inches(1.5), Inches(6), Inches(1.6)) 
    table = shape.table 
    table.cell(0, 1).text = 'Value' 
    table.cell(0, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    table.cell(1, 0).text = 'Tax liability (USD)' 
    table.cell(2, 0).text = 'Tracking error' 
    table.cell(3, 0).text = 'Total turnover' 
    
    table.cell(1, 1).text = "${:,.0f}".format(dfTradesForComponentReplacement['Tax liability'].sum()) 
    table.cell(1, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    table.cell(2, 1).text = str(round(withdrawCashTe * 100.0, 2)) + '%' 
    table.cell(2, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    table.cell(3, 1).text = "${:,.0f}".format(dfTradesForComponentReplacement['Turnover'].sum()) 
    table.cell(3, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
    
    numOfTradesRequiredForComponentReplacement = len(dfTradesForComponentReplacement.index) 
    
    shape = slide.shapes.add_table(numOfTradesRequiredForComponentReplacement + 1, 5, Inches(0.5), Inches(3.5), Inches(9), Inches(numOfTradesRequiredForComponentReplacement / 2)) 
    table = shape.table 
    
    i = 1 
    for eachCol in dfTradesForComponentReplacement.columns: 
        table.cell(0, i).text = eachCol 
        table.cell(0, i).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
        
        i = i + 1 
        
    i = 1 
    for eachIndex in dfTradesForComponentReplacement.index: 
        table.cell(i, 0).text = eachIndex 
        
        j = 1 
        for eachCol in dfTradesForComponentReplacement.columns: 
            if j == 1: 
                table.cell(i, j).text = dfTradesForComponentReplacement.loc[eachIndex, eachCol] 
            elif j == 2: 
                table.cell(i, j).text = str(round(dfTradesForComponentReplacement.loc[eachIndex, eachCol], 2)) 
            else: 
                table.cell(i, j).text = "${:,.0f}".format(dfTradesForComponentReplacement.loc[eachIndex, eachCol]) 
                
            table.cell(i, j).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT 
        
            j = j + 1 
                
        i = i + 1 
    
    todaysDate = dt.date.today().strftime('%Y%m%d') 
    
    prs.save(f'{todaysDate} Tax transition analysis - using Plaid inputs.pptx') 
